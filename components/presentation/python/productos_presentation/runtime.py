from __future__ import annotations

from copy import deepcopy
from datetime import datetime, timezone
from html import escape
import json
from pathlib import Path
import re
from typing import Any

from core.python.visual_foundations import (
    CANVAS_MAP_COMPOSITIONS as _CANVAS_MAP_COMPOSITIONS,
    MAP_COMPOSITIONS as _MAP_COMPOSITIONS,
    MATRIX_MAP_COMPOSITIONS as _MATRIX_MAP_COMPOSITIONS,
    composition_type_for_intent,
    density_mode_for_preference,
    fallback_layout_for_composition,
    html_target_profile_for_format,
    layout_variant_for_composition,
    map_layout_variant_for_composition,
    narrative_pattern_for_archetype,
    presentation_format_for_map_rendering_mode,
    presentation_mode_for_map_rendering_mode,
    presentation_theme,
    ppt_target_profile_for_format,
    density_preference_for_map_rendering_mode,
)

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:  # pragma: no cover - optional at runtime
    Presentation = None
    RGBColor = None
    MSO_CONNECTOR = None
    MSO_SHAPE = None
    PP_ALIGN = None
    Inches = None
    Pt = None


_HEADLINE_CHAR_BUDGET = {
    "hero_statement": 88,
    "decision_frame": 88,
    "risk_matrix": 92,
    "timeline_with_dependencies": 84,
    "comparison_table": 84,
    "metric_strip": 76,
    "appendix_evidence": 84,
    "summary_cards": 88,
    "evidence_grid": 88,
    "roadmap_view": 84,
    "user_journey_map": 84,
    "process_flow": 84,
    "workflow_map": 84,
    "capability_map": 84,
    "product_map": 84,
    "feature_map": 84,
    "mind_map": 84,
    "swot_matrix": 88,
    "impact_effort_matrix": 88,
}

_CORE_MESSAGE_CHAR_BUDGET = {
    "hero_statement": 140,
    "decision_frame": 128,
    "risk_matrix": 128,
    "timeline_with_dependencies": 118,
    "comparison_table": 118,
    "metric_strip": 104,
    "appendix_evidence": 110,
    "summary_cards": 118,
    "evidence_grid": 118,
    "roadmap_view": 110,
    "user_journey_map": 110,
    "process_flow": 110,
    "workflow_map": 110,
    "capability_map": 110,
    "product_map": 110,
    "feature_map": 110,
    "mind_map": 110,
    "swot_matrix": 118,
    "impact_effort_matrix": 118,
}

_ITEM_COUNT_BUDGET = {
    "hero_statement": 4,
    "decision_frame": 3,
    "risk_matrix": 4,
    "timeline_with_dependencies": 4,
    "comparison_table": 4,
    "metric_strip": 3,
    "appendix_evidence": 4,
    "summary_cards": 4,
    "evidence_grid": 4,
    "roadmap_view": 6,
    "user_journey_map": 6,
    "process_flow": 6,
    "workflow_map": 6,
    "capability_map": 6,
    "product_map": 6,
    "feature_map": 6,
    "mind_map": 6,
    "swot_matrix": 4,
    "impact_effort_matrix": 4,
}

_NUMBER_PATTERN = re.compile(r"(-?\d+(?:\.\d+)?)")


def now_iso() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z")


def theme_preset(preset: str) -> dict[str, str]:
    return presentation_theme(preset)


def _evidence_kind_for_intent(intent: str) -> str:
    return {
        "cover": "recommendation",
        "summary": "fact",
        "status": "metric",
        "timeline": "timeline_event",
        "risks": "risk",
        "decision": "decision_option",
        "portfolio": "comparison",
        "closing": "appendix_note",
    }.get(intent, "fact")


def _narrative_pattern_for_brief(presentation_brief: dict[str, Any]) -> str:
    return narrative_pattern_for_archetype(presentation_brief["presentation_archetype"])


def _confidence_for_intent(intent: str) -> str:
    return {
        "cover": "high",
        "summary": "high",
        "status": "medium",
        "timeline": "medium",
        "risks": "medium",
        "decision": "high",
        "portfolio": "medium",
        "closing": "high",
    }.get(intent, "medium")


def _story_role_for_intent(intent: str) -> str:
    return {
        "cover": "hook",
        "summary": "context",
        "status": "evidence",
        "timeline": "plan",
        "risks": "risk",
        "decision": "decision",
        "portfolio": "evidence",
        "closing": "closing",
    }.get(intent, "context")


def _composition_for_intent(intent: str) -> str:
    return composition_type_for_intent(intent)


def _density_mode(presentation_brief: dict[str, Any]) -> str:
    return density_mode_for_preference(presentation_brief["density_preference"])


def _visibility_rules(
    presentation_brief: dict[str, Any],
    outline: dict[str, Any],
) -> dict[str, Any]:
    return {
        "show_sources": True,
        "show_confidence": True,
        "show_appendix_link": bool(outline.get("appendix_link")),
        "customer_safe": presentation_brief["customer_safe"],
        "redaction_policy": presentation_brief["redaction_policy"],
    }


def _annotation_rules(outline: dict[str, Any]) -> dict[str, Any]:
    return {
        "show_confidence_badge": True,
        "show_owner_labels": outline.get("must_show_owner", False),
        "show_risk_labels": outline.get("must_show_risk", False),
        "show_source_tags": True,
    }


def _visual_tokens(
    presentation_brief: dict[str, Any],
    composition_type: str,
) -> dict[str, str]:
    theme = theme_preset(presentation_brief["theme_preset"])
    return {
        "theme_preset": theme["preset"],
        "accent": theme["accent"],
        "surface_style": "layered_panel",
        "emphasis_style": "message_first" if composition_type == "hero_statement" else "evidence_support",
    }


def _map_layout_variant(composition_type: str) -> str:
    return map_layout_variant_for_composition(composition_type)


def _layout_variant(composition_type: str, presentation_mode: str) -> str:
    return layout_variant_for_composition(composition_type, presentation_mode)


def _html_target_profile(presentation_brief: dict[str, Any], composition_type: str) -> str:
    return html_target_profile_for_format(
        presentation_brief.get("presentation_format", "both"),
        composition_type,
    )


def _ppt_target_profile(
    presentation_brief: dict[str, Any],
    composition_type: str,
    html_target_profile: str,
) -> str:
    return ppt_target_profile_for_format(
        presentation_brief.get("presentation_format", "both"),
        composition_type,
        html_target_profile,
    )


def _slide_html_target_profile(slide: dict[str, Any]) -> str:
    return slide["html_render_hints"].get("target_profile", "dual_target")


def _slide_ppt_target_profile(slide: dict[str, Any]) -> str:
    return slide["ppt_render_hints"].get("target_profile", "dual_target")


def _normalize_whitespace(text: str) -> str:
    return " ".join(str(text).split())


def _truncate_text(text: str, max_chars: int) -> str:
    normalized = _normalize_whitespace(text)
    if len(normalized) <= max_chars:
        return normalized
    clipped = normalized[: max_chars - 1].rsplit(" ", 1)[0].rstrip(" ,;:-")
    return f"{clipped}..."


def _evidence_value_budget(composition_type: str) -> int:
    if composition_type == "hero_statement":
        return 108
    if composition_type == "metric_strip":
        return 82
    if composition_type == "appendix_evidence":
        return 96
    return 92


def _compact_evidence_items(items: list[dict[str, Any]], composition_type: str) -> list[dict[str, Any]]:
    limited = []
    for item in items[: _ITEM_COUNT_BUDGET.get(composition_type, 4)]:
        limited.append(
            {
                **item,
                "label": _truncate_text(item["label"], 38),
                "value": _truncate_text(item["value"], _evidence_value_budget(composition_type)),
                "annotation": _truncate_text(item["annotation"], 58),
                "source_ref": _truncate_text(item["source_ref"], 28),
            }
        )
    return limited


def _compact_risk_items(items: list[dict[str, Any]]) -> list[dict[str, Any]]:
    limited = []
    for item in items[: _ITEM_COUNT_BUDGET["risk_matrix"]]:
        limited.append(
            {
                **item,
                "label": _truncate_text(item["label"], 78),
                "owner": _truncate_text(item["owner"], 32),
                "mitigation": _truncate_text(item["mitigation"], 92),
            }
        )
    return limited


def _compact_timeline_events(items: list[dict[str, Any]]) -> list[dict[str, Any]]:
    limited = []
    for item in items[: _ITEM_COUNT_BUDGET["timeline_with_dependencies"]]:
        limited.append(
            {
                **item,
                "label": _truncate_text(item["label"], 46),
                "dependency": _truncate_text(item["dependency"], 58),
            }
        )
    return limited


def _compact_comparison_rows(items: list[dict[str, Any]], composition_type: str) -> list[dict[str, Any]]:
    char_budget = 42 if composition_type == "comparison_table" else 32
    limited = []
    for item in items[: _ITEM_COUNT_BUDGET.get(composition_type, 4)]:
        limited.append(
            {
                **item,
                "dimension": _truncate_text(item["dimension"], 26),
                "current_state": _truncate_text(item["current_state"], char_budget),
                "target_state": _truncate_text(item["target_state"], char_budget),
                "highlight": _truncate_text(item["highlight"], 38),
            }
        )
    return limited


def _compact_map_items(
    items: list[dict[str, Any]],
    *,
    limit: int,
    label_budget: int = 28,
    description_budget: int = 88,
) -> list[dict[str, Any]]:
    compacted = []
    for item in items[:limit]:
        compacted.append(
            {
                **item,
                "label": _truncate_text(item["label"], label_budget),
                "description": _truncate_text(item.get("description", ""), description_budget) if item.get("description") else None,
            }
        )
    return compacted


def _compact_map_payload(
    map_payload: dict[str, Any],
    composition_type: str,
    presentation_mode: str,
) -> dict[str, Any]:
    limit = _ITEM_COUNT_BUDGET.get(composition_type, 6)
    payload = {
        **map_payload,
        "decision_use_case": _truncate_text(map_payload.get("decision_use_case", ""), 96),
        "certainty_notes": [_truncate_text(note, 72) for note in map_payload.get("certainty_notes", [])[:3]],
        "stages": [_truncate_text(stage, 22) for stage in map_payload.get("stages", [])[:6]],
        "lanes": [_truncate_text(lane, 22) for lane in map_payload.get("lanes", [])[:6]],
        "axes": [_truncate_text(axis, 18) for axis in map_payload.get("axes", [])[:4]],
        "items": _compact_map_items(map_payload.get("items", []), limit=limit),
        "nodes": _compact_map_items(map_payload.get("nodes", []), limit=limit),
    }
    if presentation_mode == "live":
        payload["stages"] = payload["stages"][:4]
        payload["lanes"] = payload["lanes"][:4]
        payload["items"] = payload["items"][:4]
        payload["nodes"] = payload["nodes"][:4]
    return payload


def _extract_numeric_value(text: str) -> float | None:
    match = _NUMBER_PATTERN.search(str(text))
    if not match:
        return None
    return float(match.group(1))


def _signal_index(item: dict[str, Any], confidence_state: str) -> float:
    if item.get("metric_value") is not None:
        numeric = float(item["metric_value"])
        if 0 <= numeric <= 100:
            return numeric
        return max(0.0, min(100.0, numeric))
    numeric = _extract_numeric_value(item.get("value", ""))
    if numeric is not None:
        if 0 <= numeric <= 100:
            return numeric
        return max(0.0, min(100.0, numeric))

    value = str(item.get("value", "")).lower()
    score = {"high": 76.0, "medium": 61.0, "low": 42.0}.get(confidence_state, 58.0)
    if any(token in value for token in {"improv", "ahead", "strong", "healthy", "ready"}):
        score += 12
    if any(token in value for token in {"risk", "constraint", "block", "delay", "govern"}):
        score -= 10
    return max(0.0, min(100.0, score))


def _build_chart_spec(
    chart_id: str,
    title: str,
    items: list[dict[str, Any]],
    confidence_state: str,
) -> dict[str, Any] | None:
    if not items:
        return None
    categories = [_truncate_text(item["label"], 18) for item in items[:4]]
    actual = [round(_signal_index(item, confidence_state), 1) for item in items[:4]]
    target = [min(100.0, round(value + 12, 1)) for value in actual]
    uses_heuristics = not any(_extract_numeric_value(item.get("value", "")) is not None for item in items[:4])
    return {
        "chart_id": chart_id,
        "chart_type": "bar",
        "title": title,
        "subtitle": "Derived signal index from source evidence."
        if uses_heuristics
        else "Interactive chart from numeric source values.",
        "categories": categories,
        "series": [
            {"name": "Current", "data": actual},
            {"name": "Target", "data": target},
        ],
        "interactive": True,
        "uses_heuristics": uses_heuristics,
        "unit": "index",
    }


def _source_fact_score(outline: dict[str, Any], fact: dict[str, Any]) -> int:
    score = 0
    if outline["intent"] in fact["relevance_tags"]:
        score += 2
    if fact["fact_type"] in {"metric", "status"}:
        score += 2
    if "metric" in fact["relevance_tags"]:
        score += 1
    return score


def _chart_items_from_source_material(
    presentation_brief: dict[str, Any],
    outline: dict[str, Any],
) -> list[dict[str, Any]]:
    evidence_refs = set(outline.get("evidence_refs") or presentation_brief["source_artifact_ids"])
    ranked_items = []
    for snapshot in presentation_brief.get("source_material_snapshots", []):
        if snapshot["artifact_id"] not in evidence_refs:
            continue
        for fact in snapshot["facts"]:
            if fact.get("metric_value") is None:
                continue
            score = _source_fact_score(outline, fact)
            ranked_items.append(
                (
                    score,
                    {
                        "label": fact["statement"],
                        "value": str(fact["metric_value"]),
                        "metric_value": fact["metric_value"],
                        "baseline_value": fact.get("baseline_value"),
                        "metric_target": fact.get("metric_target"),
                        "metric_unit": fact.get("metric_unit", "score"),
                    },
                )
            )
    return [item for score, item in sorted(ranked_items, key=lambda pair: pair[0], reverse=True) if score > 0][:4]


def _combo_chart_from_source_material(
    presentation_brief: dict[str, Any],
    outline: dict[str, Any],
    story_slide: dict[str, Any],
) -> dict[str, Any] | None:
    evidence_refs = set(outline.get("evidence_refs") or presentation_brief["source_artifact_ids"])
    relevant_facts = []
    combo_requested = False
    for snapshot in presentation_brief.get("source_material_snapshots", []):
        if snapshot["artifact_id"] not in evidence_refs:
            continue
        for fact in snapshot["facts"]:
            trend_points = fact.get("trend_points") or []
            if not trend_points:
                continue
            combo_requested = combo_requested or fact.get("chart_preference") == "combo" or fact.get("axis_role") == "secondary"
            relevant_facts.append((_source_fact_score(outline, fact), fact))

    if not combo_requested:
        return None
    ranked_facts = [fact for score, fact in sorted(relevant_facts, key=lambda pair: pair[0], reverse=True) if score > 0][:3]
    if not ranked_facts:
        return None
    categories = [point["label"] for point in ranked_facts[0]["trend_points"]]
    series = []
    primary_axis = False
    secondary_axis = False
    for fact in ranked_facts:
        if [point["label"] for point in fact["trend_points"]] != categories:
            continue
        y_axis_index = 1 if fact.get("axis_role") == "secondary" else 0
        primary_axis = primary_axis or y_axis_index == 0
        secondary_axis = secondary_axis or y_axis_index == 1
        series.append(
            {
                "name": _truncate_text(fact["statement"], 20),
                "type": "bar" if y_axis_index == 1 else "line",
                "y_axis_index": y_axis_index,
                "data": [float(point["value"]) for point in fact["trend_points"]],
            }
        )
    if not series or not (primary_axis and secondary_axis):
        return None
    return {
        "chart_id": f"{story_slide['slide_id']}_combo_chart",
        "chart_type": "combo",
        "title": "Operating trend and load",
        "subtitle": "Dual-axis trend from numeric source values.",
        "categories": [_truncate_text(label, 14) for label in categories],
        "series": series,
        "interactive": True,
        "uses_heuristics": False,
        "unit": "mixed",
        "y_axes": ["score", "count"],
    }


def _trend_chart_from_source_material(
    presentation_brief: dict[str, Any],
    outline: dict[str, Any],
    story_slide: dict[str, Any],
) -> dict[str, Any] | None:
    evidence_refs = set(outline.get("evidence_refs") or presentation_brief["source_artifact_ids"])
    trend_facts = []
    for snapshot in presentation_brief.get("source_material_snapshots", []):
        if snapshot["artifact_id"] not in evidence_refs:
            continue
        for fact in snapshot["facts"]:
            trend_points = fact.get("trend_points") or []
            if not trend_points:
                continue
            if fact.get("axis_role") == "secondary" or fact.get("chart_preference") == "combo":
                continue
            score = _source_fact_score(outline, fact)
            trend_facts.append((score, fact))

    ranked_facts = [fact for score, fact in sorted(trend_facts, key=lambda pair: pair[0], reverse=True) if score > 0][:2]
    if not ranked_facts:
        return None

    categories = [point["label"] for point in ranked_facts[0]["trend_points"]]
    series = []
    target_series = []
    baseline_series = []
    for fact in ranked_facts:
        if [point["label"] for point in fact["trend_points"]] != categories:
            continue
        series.append(
            {
                "name": _truncate_text(fact["statement"], 20),
                "type": "line",
                "data": [float(point["value"]) for point in fact["trend_points"]],
            }
        )
        if fact.get("metric_target") is not None:
            target_series.append(float(fact["metric_target"]))
        if fact.get("baseline_value") is not None:
            baseline_series.append(float(fact["baseline_value"]))
    if not series:
        return None
    overlay_series = []
    if target_series:
        overlay_series.append(
            {
                "name": "Target",
                "type": "line",
                "data": [round(sum(target_series) / len(target_series), 1)] * len(categories),
            }
        )
    if baseline_series:
        overlay_series.append(
            {
                "name": "Baseline",
                "type": "line",
                "data": [round(sum(baseline_series) / len(baseline_series), 1)] * len(categories),
            }
        )

    return {
        "chart_id": f"{story_slide['slide_id']}_trend_chart",
        "chart_type": "line",
        "title": "Operating trend",
        "subtitle": "Interactive trend from numeric source values.",
        "categories": [_truncate_text(label, 14) for label in categories],
        "series": series + overlay_series,
        "interactive": True,
        "uses_heuristics": False,
        "unit": "score",
    }


def _stacked_chart_from_source_material(
    presentation_brief: dict[str, Any],
    outline: dict[str, Any],
    story_slide: dict[str, Any],
) -> dict[str, Any] | None:
    evidence_refs = set(outline.get("evidence_refs") or presentation_brief["source_artifact_ids"])
    stacked_facts = []
    for snapshot in presentation_brief.get("source_material_snapshots", []):
        if snapshot["artifact_id"] not in evidence_refs:
            continue
        for fact in snapshot["facts"]:
            components = fact.get("metric_components") or []
            if not components:
                continue
            if fact.get("chart_preference") not in {None, "stacked_bar", "comparison"}:
                continue
            stacked_facts.append((_source_fact_score(outline, fact), fact))
    ranked_facts = [
        fact
        for score, fact in sorted(
            stacked_facts,
            key=lambda pair: (
                0 if "current" in pair[1]["statement"].lower() else 1 if "target" in pair[1]["statement"].lower() else 2,
                0 if pair[1]["fact_type"] == "status" else 1,
                -pair[0],
                pair[1]["statement"],
            ),
        )
        if score > 0
    ][:3]
    if not ranked_facts:
        return None
    categories = [_truncate_text(fact["statement"], 16) for fact in ranked_facts]
    component_labels = []
    for fact in ranked_facts:
        for component in fact["metric_components"]:
            if component["label"] not in component_labels:
                component_labels.append(component["label"])
    series = []
    for label in component_labels[:4]:
        data = []
        for fact in ranked_facts:
            component_lookup = {component["label"]: float(component["value"]) for component in fact["metric_components"]}
            data.append(component_lookup.get(label, 0.0))
        series.append({"name": _truncate_text(label, 16), "type": "bar", "stack": "total", "data": data})
    return {
        "chart_id": f"{story_slide['slide_id']}_stacked_chart",
        "chart_type": "stacked_bar",
        "title": "Dependency composition",
        "subtitle": "Stacked comparison of current and target states.",
        "categories": categories,
        "series": series,
        "interactive": True,
        "uses_heuristics": False,
        "unit": "count",
    }


def _heatmap_chart_from_source_material(
    presentation_brief: dict[str, Any],
    outline: dict[str, Any],
    story_slide: dict[str, Any],
) -> dict[str, Any] | None:
    evidence_refs = set(outline.get("evidence_refs") or presentation_brief["source_artifact_ids"])
    cells = []
    for snapshot in presentation_brief.get("source_material_snapshots", []):
        if snapshot["artifact_id"] not in evidence_refs:
            continue
        for fact in snapshot["facts"]:
            for cell in fact.get("portfolio_cells") or []:
                cells.append(cell)
    if not cells:
        return None
    x_categories = list(dict.fromkeys(cell["x_label"] for cell in cells))
    y_categories = list(dict.fromkeys(cell["y_label"] for cell in cells))
    points = []
    for cell in cells:
        points.append(
            {
                "x": x_categories.index(cell["x_label"]),
                "y": y_categories.index(cell["y_label"]),
                "value": float(cell["value"]),
                "label": cell.get("label", ""),
            }
        )
    return {
        "chart_id": f"{story_slide['slide_id']}_heatmap_chart",
        "chart_type": "heatmap",
        "title": "Portfolio heatmap",
        "subtitle": "Interactive portfolio intensity grid.",
        "categories": [],
        "x_categories": [_truncate_text(label, 14) for label in x_categories],
        "y_categories": [_truncate_text(label, 14) for label in y_categories],
        "series": [],
        "heatmap_points": points,
        "interactive": True,
        "uses_heuristics": False,
        "unit": "score",
    }


def _build_chart_spec_for_slide(
    presentation_brief: dict[str, Any],
    outline: dict[str, Any],
    story_slide: dict[str, Any],
    composition_type: str,
    payload: dict[str, Any],
) -> dict[str, Any] | None:
    if composition_type not in {"metric_strip", "comparison_table"}:
        return None
    if composition_type == "comparison_table":
        heatmap_chart = _heatmap_chart_from_source_material(presentation_brief, outline, story_slide)
        if heatmap_chart:
            return heatmap_chart
        stacked_chart = _stacked_chart_from_source_material(presentation_brief, outline, story_slide)
        if stacked_chart:
            return stacked_chart
    combo_chart = _combo_chart_from_source_material(presentation_brief, outline, story_slide)
    if combo_chart:
        return combo_chart
    trend_chart = _trend_chart_from_source_material(presentation_brief, outline, story_slide)
    if trend_chart:
        return trend_chart
    source_metric_items = _chart_items_from_source_material(presentation_brief, outline)
    if source_metric_items:
        return {
            "chart_id": f"{story_slide['slide_id']}_chart",
            "chart_type": "bar",
            "title": "Operating metrics",
            "subtitle": "Interactive chart from numeric source values.",
            "categories": [_truncate_text(item["label"], 18) for item in source_metric_items],
            "series": [
                {
                    "name": "Current",
                    "type": "bar",
                    "data": [round(float(item["metric_value"]), 1) for item in source_metric_items],
                }
            ]
            + (
                [
                    {
                        "name": "Baseline",
                        "type": "line",
                        "data": [round(float(item["baseline_value"]), 1) for item in source_metric_items],
                    }
                ]
                if all(item.get("baseline_value") is not None for item in source_metric_items)
                else []
            )
            + (
                [
                    {
                        "name": "Target",
                        "type": "line",
                        "data": [round(float(item["metric_target"]), 1) for item in source_metric_items],
                    }
                ]
                if all(item.get("metric_target") is not None for item in source_metric_items)
                else []
            ),
            "interactive": True,
            "uses_heuristics": False,
            "unit": source_metric_items[0].get("metric_unit", "score"),
        }
    if composition_type == "metric_strip":
        return payload.get("chart_spec")
    return None


def _apply_render_budgets(slide: dict[str, Any], presentation_brief: dict[str, Any]) -> dict[str, Any]:
    composition_type = slide["composition_type"]
    payload = dict(slide["composition_payload"])

    fitted_slide = {
        **slide,
        "headline": _truncate_text(slide["headline"], _HEADLINE_CHAR_BUDGET.get(composition_type, 88)),
        "core_message": _truncate_text(slide["core_message"], _CORE_MESSAGE_CHAR_BUDGET.get(composition_type, 120)),
        "speaker_notes": _truncate_text(slide["speaker_notes"], 220),
        "visual_direction": _truncate_text(slide["visual_direction"], 100),
        "source_refs": [_truncate_text(source, 28) for source in slide["source_refs"][:4]],
    }

    payload["primary_claim"] = _truncate_text(payload.get("primary_claim", fitted_slide["core_message"]), 132)
    payload["summary"] = _truncate_text(payload.get("summary", fitted_slide["core_message"]), 110)
    if payload.get("recommendation"):
        payload["recommendation"] = _truncate_text(payload["recommendation"], 74)
    if payload.get("decision_ask"):
        payload["decision_ask"] = _truncate_text(payload["decision_ask"], 68)
    if payload.get("options"):
        payload["options"] = [
            {
                **option,
                "label": _truncate_text(option["label"], 22),
                "summary": _truncate_text(option["summary"], 76),
                "tradeoffs": [_truncate_text(tradeoff, 44) for tradeoff in option["tradeoffs"][:3]],
            }
            for option in payload["options"][:2]
        ]

    payload["evidence_items"] = _compact_evidence_items(payload.get("evidence_items", []), composition_type)
    payload["risk_items"] = _compact_risk_items(payload.get("risk_items", []))
    payload["timeline_events"] = _compact_timeline_events(payload.get("timeline_events", []))
    payload["comparison_rows"] = _compact_comparison_rows(payload.get("comparison_rows", []), composition_type)
    if payload.get("map_payload"):
        payload["map_payload"] = _compact_map_payload(
            payload["map_payload"],
            composition_type,
            presentation_brief["presentation_mode"],
        )
    if payload.get("chart_spec"):
        payload["chart_spec"] = {
            **payload["chart_spec"],
            "categories": [_truncate_text(category, 18) for category in payload["chart_spec"]["categories"][:4]],
            "x_categories": [_truncate_text(category, 14) for category in payload["chart_spec"].get("x_categories", [])[:6]],
            "y_categories": [_truncate_text(category, 14) for category in payload["chart_spec"].get("y_categories", [])[:6]],
            "series": [
                {
                    **series,
                    "name": _truncate_text(series["name"], 18),
                    "data": series["data"][:4],
                }
                for series in payload["chart_spec"]["series"][:4]
            ],
            "heatmap_points": payload["chart_spec"].get("heatmap_points", [])[:16],
            "y_axes": payload["chart_spec"].get("y_axes", [])[:2],
        }

    if presentation_brief["presentation_mode"] == "live":
        payload["evidence_items"] = payload["evidence_items"][:2]
        payload["risk_items"] = payload["risk_items"][:2]
        payload["timeline_events"] = payload["timeline_events"][:3]
        payload["comparison_rows"] = payload["comparison_rows"][:3]
    elif presentation_brief["presentation_mode"] in {"memo", "one_page", "meeting_brief"}:
        fitted_slide["core_message"] = _truncate_text(fitted_slide["core_message"], 104)
        payload["evidence_items"] = payload["evidence_items"][:3]

    fitted_slide["composition_payload"] = payload
    fitted_slide["html_render_hints"] = {
        **slide["html_render_hints"],
        "max_visible_evidence": len(payload["evidence_items"]),
    }
    fitted_slide["ppt_render_hints"] = {
        **slide["ppt_render_hints"],
        "max_visible_evidence": min(len(payload["evidence_items"]), 4),
    }
    return fitted_slide


def _build_composition_payload(
    outline: dict[str, Any],
    story_slide: dict[str, Any],
) -> dict[str, Any]:
    evidence_items = []
    for index, point in enumerate(story_slide["supporting_points"]):
        source_refs = story_slide["evidence_refs"] or ["source_unavailable"]
        evidence_items.append(
            {
                "evidence_id": f"{story_slide['slide_id']}_evidence_{index + 1}",
                "label": outline["proof_requirements"][min(index, len(outline["proof_requirements"]) - 1)]
                if outline.get("proof_requirements")
                else f"Support point {index + 1}",
                "evidence_type": _evidence_kind_for_intent(outline["intent"]),
                "value": point,
                "source_ref": source_refs[min(index, len(source_refs) - 1)],
                "annotation": outline.get("audience_question") or story_slide["claim"],
            }
        )

    if story_slide["composition_type"] == "decision_frame":
        options = [
            {
                "option_id": f"{story_slide['slide_id']}_recommended",
                "label": "Recommended path",
                "stance": "recommended",
                "summary": story_slide["core_message"],
                "tradeoffs": outline.get("key_points") or story_slide["supporting_points"],
            }
        ]
        presentation_brief_option = outline.get("proof_requirements")
        if presentation_brief_option:
            options.append(
                {
                    "option_id": f"{story_slide['slide_id']}_alternative",
                    "label": "Alternative",
                    "stance": "alternative",
                    "summary": f"Delay or narrow the ask until {presentation_brief_option[0].lower()}.",
                    "tradeoffs": [story_slide["supporting_points"][0], "Reduces immediate scope but slows decision velocity."],
                }
            )
        return {
            "primary_claim": story_slide["claim"],
            "summary": story_slide["core_message"],
            "recommendation": story_slide["core_message"],
            "options": options,
            "decision_ask": story_slide["cta"],
            "evidence_items": evidence_items,
        }
    if story_slide["composition_type"] == "risk_matrix":
        return {
            "primary_claim": story_slide["claim"],
            "summary": story_slide["core_message"],
            "risk_items": [
                {
                    "risk_id": f"{story_slide['slide_id']}_risk_{index + 1}",
                    "label": point,
                    "likelihood": "medium",
                    "impact": "high",
                    "owner": "PM / presentation-system owner" if outline.get("must_show_owner") else "owner_unassigned",
                    "mitigation": story_slide["proof_strategy"],
                }
                for index, point in enumerate(story_slide["supporting_points"])
            ],
            "evidence_items": evidence_items,
        }
    if story_slide["composition_type"] == "timeline_with_dependencies":
        return {
            "primary_claim": story_slide["claim"],
            "summary": story_slide["core_message"],
            "timeline_events": [
                {
                    "event_id": f"{story_slide['slide_id']}_event_{index + 1}",
                    "label": point,
                    "timing": f"T{index + 1}",
                    "dependency": "Narrative planning and composition contracts",
                    "confidence_state": story_slide["confidence_state"],
                }
                for index, point in enumerate(story_slide["supporting_points"])
            ],
            "evidence_items": evidence_items,
        }
    if story_slide["composition_type"] == "comparison_table":
        return {
            "primary_claim": story_slide["claim"],
            "summary": story_slide["core_message"],
            "comparison_rows": [
                {
                    "dimension": outline["proof_requirements"][min(index, len(outline["proof_requirements"]) - 1)]
                    if outline.get("proof_requirements")
                    else f"Dimension {index + 1}",
                    "current_state": point,
                    "target_state": story_slide["claim"],
                    "highlight": "Gap to close now",
                }
                for index, point in enumerate(story_slide["supporting_points"])
            ],
            "evidence_items": evidence_items,
        }
    if story_slide["composition_type"] == "metric_strip":
        chart_items = evidence_items[:3]
        return {
            "primary_claim": story_slide["claim"],
            "summary": story_slide["core_message"],
            "evidence_items": evidence_items,
            "comparison_rows": [
                {
                    "dimension": item["label"],
                    "current_state": item["value"],
                    "target_state": story_slide["core_message"],
                    "highlight": item["annotation"],
                }
                for item in evidence_items[:3]
            ],
            "chart_spec": _build_chart_spec(
                f"{story_slide['slide_id']}_chart",
                "Operating signal index",
                chart_items,
                story_slide["confidence_state"],
            ),
        }
    if story_slide["composition_type"] == "appendix_evidence":
        return {
            "primary_claim": story_slide["claim"],
            "summary": story_slide["core_message"],
            "evidence_items": evidence_items,
            "comparison_rows": [
                {
                    "dimension": "Appendix trace",
                    "current_state": item["source_ref"],
                    "target_state": item["value"],
                    "highlight": item["annotation"],
                }
                for item in evidence_items[:3]
            ],
        }
    return {
        "primary_claim": story_slide["claim"],
        "evidence_items": evidence_items,
        "summary": story_slide["core_message"],
    }


def _payload_evidence_items(slide: dict[str, Any]) -> list[dict[str, Any]]:
    payload = slide.get("composition_payload", {})
    if "evidence_items" in payload:
        return payload["evidence_items"]
    if "risk_items" in payload:
        return [
            {
                "evidence_id": item["risk_id"],
                "label": item["label"],
                "evidence_type": "risk",
                "value": item["mitigation"],
                "source_ref": slide["source_refs"][0] if slide["source_refs"] else "source_unavailable",
                "annotation": f"{item['likelihood']} likelihood / {item['impact']} impact",
            }
            for item in payload["risk_items"]
        ]
    return []


def _clamp_score(value: float) -> float:
    return max(0.0, min(5.0, round(value, 1)))


def _narrative_quality(render_spec: dict[str, Any], presentation_brief: dict[str, Any]) -> tuple[float, list[str]]:
    notes = []
    slides = render_spec["slides"]
    score = 2.5
    if slides:
        first_slide = slides[0]
        if presentation_brief["decision_required"].lower()[:20] in first_slide["core_message"].lower() or "recommend" in first_slide["core_message"].lower():
            score += 1.0
        else:
            notes.append("The opening slide does not surface the recommendation or decision ask strongly enough.")
        if first_slide["composition_payload"].get("primary_claim"):
            score += 0.4
    if all(slide["composition_payload"].get("primary_claim") for slide in slides):
        score += 0.4
    if any(slide["composition_type"] == "risk_matrix" for slide in slides):
        score += 0.4
    else:
        notes.append("No explicit risk-treatment slide is present.")
    if presentation_brief.get("required_objections"):
        score += 0.3
    else:
        notes.append("The brief does not specify objections to address.")
    return _clamp_score(score), notes


def _evidence_quality(render_spec: dict[str, Any]) -> tuple[float, list[str]]:
    notes = []
    score = 2.0
    evidence_counts = [len(_payload_evidence_items(slide)) for slide in render_spec["slides"]]
    if evidence_counts and min(evidence_counts) >= 1:
        score += 1.0
    else:
        notes.append("At least one slide is missing visible evidence items.")
    if all(slide["source_refs"] for slide in render_spec["slides"]):
        score += 1.0
    else:
        notes.append("At least one slide is missing source references.")
    if any(slide["confidence_state"] in {"medium", "low"} for slide in render_spec["slides"]):
        score += 0.4
    else:
        notes.append("The deck does not show any meaningful confidence variation.")
    if any(slide["composition_type"] == "risk_matrix" for slide in render_spec["slides"]):
        score += 0.4
    return _clamp_score(score), notes


def _audience_fit_quality(render_spec: dict[str, Any], presentation_brief: dict[str, Any]) -> tuple[float, list[str]]:
    notes = []
    score = 2.5
    density_mode = _density_mode(presentation_brief)
    if all(slide["density_mode"] == density_mode for slide in render_spec["slides"]):
        score += 0.7
    else:
        notes.append("Slides are not using a consistent density mode.")
    if presentation_brief["presenter_mode"] == "self_explanatory" and presentation_brief["presentation_mode"] == "async":
        score += 0.7
    if presentation_brief["customer_safe"] and presentation_brief["redaction_policy"] == "none":
        notes.append("Customer-safe outputs require a non-none redaction policy.")
        score -= 1.2
    else:
        score += 0.4
    if presentation_brief["audience_type"] == "leadership":
        if any(slide["composition_type"] in {"decision_frame", "risk_matrix", "hero_statement"} for slide in render_spec["slides"]):
            score += 0.5
        else:
            notes.append("Leadership deck is missing high-signal decision or risk structures.")
    return _clamp_score(score), notes


def _visual_clarity_quality(render_spec: dict[str, Any]) -> tuple[float, list[str]]:
    notes = []
    score = 2.5
    composition_types = {slide["composition_type"] for slide in render_spec["slides"]}
    if len(composition_types) >= min(2, len(render_spec["slides"])):
        score += 0.9
    else:
        notes.append("The deck uses too few distinct composition types.")
    if all(slide["layout_variant"] for slide in render_spec["slides"]):
        score += 0.5
    if any(slide["annotation_rules"].get("show_source_tags") for slide in render_spec["slides"]):
        score += 0.4
    if any(slide["annotation_rules"].get("show_risk_labels") for slide in render_spec["slides"]):
        score += 0.3
    visual_payload_slides = sum(1 for slide in render_spec["slides"] if _slide_has_visual_payload(slide))
    if visual_payload_slides >= max(1, len(render_spec["slides"]) - 1):
        score += 0.6
    else:
        notes.append("Several slides still rely on text containers instead of explicit visual primitives.")
    return _clamp_score(score), notes


def _slide_has_visual_payload(slide: dict[str, Any]) -> bool:
    payload = slide.get("composition_payload", {})
    if slide["composition_type"] == "hero_statement":
        return bool(payload.get("recommendation") or payload.get("evidence_items"))
    if slide["composition_type"] == "risk_matrix":
        return bool(payload.get("risk_items"))
    if slide["composition_type"] == "timeline_with_dependencies":
        return bool(payload.get("timeline_events"))
    if slide["composition_type"] in {"comparison_table", "metric_strip", "appendix_evidence"}:
        return bool(payload.get("comparison_rows"))
    if slide["composition_type"] in _MAP_COMPOSITIONS:
        map_payload = payload.get("map_payload") or {}
        return bool(
            map_payload.get("items")
            or map_payload.get("nodes")
            or map_payload.get("stages")
            or map_payload.get("lanes")
            or map_payload.get("axes")
        )
    return bool(payload.get("evidence_items"))


def _normalize_source_material(presentation_brief: dict[str, Any]) -> list[dict[str, Any]]:
    normalized_facts = []
    for snapshot in presentation_brief.get("source_material_snapshots", []):
        for fact in snapshot["facts"]:
            normalized_facts.append(
                {
                    "artifact_id": snapshot["artifact_id"],
                    "artifact_type": snapshot["artifact_type"],
                    "fact_id": fact["fact_id"],
                    "fact_type": fact["fact_type"],
                    "statement": fact["statement"],
                    "claim_mode": fact.get("claim_mode", "observed"),
                    "validation_note": fact.get("validation_note", ""),
                    "metric_value": fact.get("metric_value"),
                    "relevance_tags": fact["relevance_tags"],
                }
            )
    return normalized_facts


def _detect_contradictions(normalized_facts: list[dict[str, Any]]) -> list[dict[str, Any]]:
    contradictions = []
    grouped: dict[str, set[str]] = {}
    for fact in normalized_facts:
        for tag in fact["relevance_tags"]:
            grouped.setdefault(tag, set()).add(fact["fact_type"])
    for tag, fact_types in grouped.items():
        if "risk" in fact_types and "decision" in fact_types:
            contradictions.append(
                {
                    "summary": f"Facts tagged '{tag}' contain both decision and risk signals that may need explicit reconciliation.",
                    "severity": "medium",
                }
            )
        if "status" in fact_types and "constraint" in fact_types:
            contradictions.append(
                {
                    "summary": f"Facts tagged '{tag}' mix status and constraint signals; presentation should avoid overstating certainty.",
                    "severity": "low",
                }
            )
    return contradictions


def _explicit_contradictions_from_brief(presentation_brief: dict[str, Any]) -> list[dict[str, Any]]:
    return [
        {
            "summary": summary,
            "severity": "high",
        }
        for summary in presentation_brief.get("contradiction_summaries", [])
    ]


def _fallback_normalized_facts(presentation_brief: dict[str, Any]) -> list[dict[str, Any]]:
    normalized_facts = []
    for outline in presentation_brief["slide_outlines"]:
        for index, point in enumerate(outline.get("key_points") or [presentation_brief["objective"]]):
            normalized_facts.append(
                {
                    "artifact_id": (outline.get("evidence_refs") or presentation_brief["source_artifact_ids"])[0],
                    "artifact_type": "presentation_outline",
                    "fact_id": f"{outline['slide_id']}_fallback_{index + 1}",
                    "fact_type": _evidence_kind_for_intent(outline["intent"]),
                    "statement": point,
                    "metric_value": None,
                    "relevance_tags": [outline["intent"], "fallback"],
                }
            )
    return normalized_facts


def _select_source_facts(
    presentation_brief: dict[str, Any],
    outline: dict[str, Any],
    normalized_facts: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    evidence_refs = set(outline.get("evidence_refs") or presentation_brief["source_artifact_ids"])
    intent = outline["intent"]
    proof_terms = {term.lower() for term in outline["proof_requirements"]}
    claim_terms = {term.lower() for term in outline["claim"].split()}
    scored_facts = []

    for fact in normalized_facts:
        score = 0
        if fact["artifact_id"] in evidence_refs:
            score += 3
        if intent in fact["relevance_tags"]:
            score += 2
        if any(tag in {"evidence", "recommendation", "governance", "status", "risk"} for tag in fact["relevance_tags"]):
            score += 1
        statement_lower = fact["statement"].lower()
        if any(term in statement_lower for term in proof_terms):
            score += 2
        if any(term in statement_lower for term in claim_terms if len(term) > 4):
            score += 1
        if presentation_brief["customer_safe"] and fact["fact_type"] in {"risk", "constraint"}:
            score -= 1
        if outline["intent"] == "closing" and fact["fact_type"] in {"decision", "timeline", "observation"}:
            score += 1
        if outline["intent"] == "status" and fact["fact_type"] in {"status", "metric"}:
            score += 2
        scored_facts.append((score, fact))

    ranked = [fact for score, fact in sorted(scored_facts, key=lambda item: item[0], reverse=True) if score > 0]
    if ranked:
        return ranked[: max(2, min(4, len(ranked)))]
    return [fact for fact in normalized_facts[:2]]


def _derive_confidence_from_facts(facts: list[dict[str, Any]], intent: str) -> str:
    if not facts:
        return _confidence_for_intent(intent)
    fact_types = {fact["fact_type"] for fact in facts}
    claim_modes = {fact.get("claim_mode", "observed") for fact in facts}
    if "hypothesis" in claim_modes:
        return "low"
    if "inferred" in claim_modes:
        return "medium"
    if "risk" in fact_types or "constraint" in fact_types:
        return "medium"
    if len(facts) >= 2:
        return "high"
    return "medium"


def _story_arc_for_brief(presentation_brief: dict[str, Any]) -> list[str]:
    pattern = _narrative_pattern_for_brief(presentation_brief)
    if pattern == "answer_first":
        return [
            "Start with the recommendation and decision ask.",
            "Prove the recommendation with grounded operating facts.",
            "Close with risks, boundaries, and the controlled path forward.",
        ]
    if pattern == "option_comparison":
        return [
            "Frame the operating question and comparison baseline.",
            "Compare options or states using evidence-led contrasts.",
            "Close with the chosen path and the tradeoffs accepted.",
        ]
    return [
        "Establish the current situation.",
        "Show the complication or pressure point using source facts.",
        "Answer with the recommended path and bounded next steps.",
    ]


def _synthetic_outline_for_story_slide(
    presentation_brief: dict[str, Any],
    story_slide: dict[str, Any],
) -> dict[str, Any]:
    proof_requirements = [item.strip() for item in story_slide.get("proof_strategy", "").split(";") if item.strip()]
    if not proof_requirements:
        proof_requirements = ["Keep explicit source labels and confidence visible."]
    return {
        "slide_id": story_slide["slide_id"],
        "title": story_slide["title"],
        "intent": "risks" if story_slide.get("story_role") in {"risk", "appendix"} else "summary",
        "audience_question": story_slide["audience_question_answered"],
        "claim": story_slide["claim"],
        "proof_requirements": proof_requirements,
        "core_message": story_slide["core_message"],
        "key_points": story_slide["supporting_points"],
        "evidence_refs": story_slide["evidence_refs"],
        "speaker_notes": story_slide["speaker_notes"],
        "visual_direction": story_slide["visual_direction"],
        "cta": story_slide["cta"],
        "appendix_link": story_slide["appendix_reference"],
    }


def _transition_text(previous_title: str | None, current_title: str, pattern: str) -> str:
    if previous_title is None:
        return "Open with the answer first." if pattern == "answer_first" else "Open by framing the situation before adding detail."
    if pattern == "option_comparison":
        return f"Move from {previous_title} to {current_title} by tightening the tradeoff discussion."
    if pattern == "scq":
        return f"Move from {previous_title} to {current_title} by increasing the complication and then resolving it."
    return f"Build from {previous_title} into {current_title} without breaking the recommendation thread."


def _claim_posture_summary(facts: list[dict[str, Any]]) -> str:
    if not facts:
        return "Claim posture: observed facts only."
    claim_modes = [fact.get("claim_mode", "observed") for fact in facts]
    if "hypothesis" in claim_modes:
        return "Claim posture: includes hypothesis-level inputs, so keep uncertainty explicit."
    if "inferred" in claim_modes:
        return "Claim posture: includes inferred inputs that still need external or PM validation."
    return "Claim posture: observed facts support the current narrative."


def _validation_note_summary(facts: list[dict[str, Any]]) -> str:
    notes = [fact.get("validation_note", "").strip() for fact in facts if fact.get("validation_note", "").strip()]
    if not notes:
        return ""
    return f"Validation next: {' '.join(notes[:2])}"


def build_evidence_pack(presentation_brief: dict[str, Any]) -> dict[str, Any]:
    normalized_facts = _normalize_source_material(presentation_brief) or _fallback_normalized_facts(presentation_brief)
    contradictions = _detect_contradictions(normalized_facts)
    contradictions.extend(_explicit_contradictions_from_brief(presentation_brief))
    seen_contradictions = set()
    contradictions = [
        item
        for item in contradictions
        if not (item["summary"] in seen_contradictions or seen_contradictions.add(item["summary"]))
    ]
    evidence_units = []
    priority_evidence = []
    deferred_evidence = []
    risk_items = []
    timeline_events = []
    decision_options = []

    for outline in presentation_brief["slide_outlines"]:
        selected_facts = _select_source_facts(presentation_brief, outline, normalized_facts)
        supporting_points = [fact["statement"] for fact in selected_facts] or (outline.get("key_points") or [presentation_brief["objective"]])
        source_refs = outline.get("evidence_refs") or presentation_brief["source_artifact_ids"]
        confidence_state = _derive_confidence_from_facts(selected_facts, outline["intent"])
        claim_posture_summary = _claim_posture_summary(selected_facts)
        validation_note_summary = _validation_note_summary(selected_facts)
        for index, point in enumerate(supporting_points):
            selected_fact = selected_facts[min(index, len(selected_facts) - 1)] if selected_facts else None
            evidence_unit = {
                "evidence_id": f"{outline['slide_id']}_evidence_{index + 1}",
                "slide_id": outline["slide_id"],
                "claim": outline["claim"],
                "evidence_type": _evidence_kind_for_intent(outline["intent"]),
                "value": point,
                "source_ref": selected_fact["artifact_id"] if selected_fact else source_refs[min(index, len(source_refs) - 1)],
                "confidence_state": confidence_state,
                "confidence_reason": (
                    f"Confidence is derived from the number, type, and claim posture of normalized source facts available for this claim. {claim_posture_summary}"
                ),
                "claim_mode": selected_fact.get("claim_mode", "observed") if selected_fact else "observed",
                "freshness": "current_quarter" if selected_fact else "unspecified",
                "materiality": "primary" if index == 0 else "supporting",
                "proof_requirement": outline["proof_requirements"][min(index, len(outline["proof_requirements"]) - 1)],
                "counterevidence": [],
            }
            if selected_fact and selected_fact.get("validation_note", "").strip():
                evidence_unit["validation_note"] = selected_fact["validation_note"].strip()
            evidence_units.append(evidence_unit)
        if validation_note_summary:
            deferred_evidence.append(
                {
                    "slide_id": outline["slide_id"],
                    "reason": "Validation notes that should remain visible in appendix or speaker guidance.",
                    "deferred_points": [validation_note_summary],
                }
            )
        priority_evidence.append(
            {
                "slide_id": outline["slide_id"],
                "claim": outline["claim"],
                "proof_points": supporting_points[:2],
            }
        )
        deferred_evidence.append(
            {
                "slide_id": outline["slide_id"],
                "reason": "Appendix candidate or deeper proof that should not dilute the primary narrative.",
                "deferred_points": supporting_points[2:] or [f"Expanded source trace for {outline['slide_id']}"],
            }
        )
        if outline["intent"] == "risks":
            risk_items.extend(
                {
                    "risk_id": f"{outline['slide_id']}_risk_{index + 1}",
                    "label": point,
                    "impact": "high",
                    "likelihood": "medium",
                    "owner": "PM / presentation-system owner" if outline.get("must_show_owner") else "owner_unassigned",
                    "source_ref": source_refs[min(index, len(source_refs) - 1)],
                }
                for index, point in enumerate(supporting_points)
            )
        if outline["intent"] == "timeline":
            timeline_events.extend(
                {
                    "event_id": f"{outline['slide_id']}_event_{index + 1}",
                    "label": point,
                    "timing": f"T{index + 1}",
                    "dependency": "Narrative planning and render-spec upgrades",
                    "source_ref": source_refs[min(index, len(source_refs) - 1)],
                }
                for index, point in enumerate(supporting_points)
            )
        if outline["intent"] == "decision":
            decision_options.append(
                {
                    "option_id": f"{outline['slide_id']}_recommended",
                    "label": "Recommended path",
                    "stance": "recommended",
                    "summary": outline.get("core_message") or outline["claim"],
                    "source_refs": source_refs,
                }
            )

    gaps = []
    if presentation_brief["presentation_mode"] in {"one_page", "meeting_brief", "customer_narrative"}:
        gaps.append("Condensed or audience-safe modes require stronger evidence prioritization and redaction discipline.")
    gaps.extend(presentation_brief.get("known_gaps", []))
    if presentation_brief.get("external_research_questions"):
        gaps.append(
            "Open external research questions remain: "
            + "; ".join(presentation_brief["external_research_questions"][:2])
        )
    if presentation_brief.get("contradiction_summaries"):
        gaps.append(
            "Conflicted external evidence remains visible: "
            + "; ".join(presentation_brief["contradiction_summaries"][:2])
        )

    return {
        "schema_version": "1.0.0",
        "evidence_pack_id": f"evidence_pack_{presentation_brief['presentation_brief_id']}",
        "workspace_id": presentation_brief["workspace_id"],
        "presentation_brief_id": presentation_brief["presentation_brief_id"],
        "source_artifact_ids": presentation_brief["source_artifact_ids"],
        "evidence_units": evidence_units,
        "priority_evidence": priority_evidence,
        "deferred_evidence": deferred_evidence,
        "metrics": [{"metric_name": "slide_count", "value": len(presentation_brief["slide_outlines"]), "unit": "slides"}],
        "confidence_flags": [
            "Publishing output must preserve explicit source references.",
            "Confidence is now source-fact aware, but contradiction detection still needs deeper artifact-level reasoning.",
            *(
                ["Some presentation claims still depend on inferred or hypothesis-level inputs and should keep uncertainty visible."]
                if any(fact.get("claim_mode") in {"inferred", "hypothesis"} for fact in normalized_facts)
                else []
            ),
            *(
                ["Conflicted external evidence should be made explicit in the narrative and decision framing."]
                if contradictions
                else []
            ),
        ],
        "contradictions": contradictions,
        "gaps": gaps,
        "risk_items": risk_items,
        "timeline_events": timeline_events,
        "decision_options": decision_options,
        "generated_at": now_iso(),
    }


def build_presentation_story(
    presentation_brief: dict[str, Any],
    evidence_pack: dict[str, Any],
) -> dict[str, Any]:
    evidence_lookup: dict[str, list[dict[str, Any]]] = {}
    for unit in evidence_pack["evidence_units"]:
        evidence_lookup.setdefault(unit["slide_id"], []).append(unit)
    story_slides = []
    narrative_pattern = _narrative_pattern_for_brief(presentation_brief)
    contradiction_summaries = presentation_brief.get("contradiction_summaries", [])

    for outline in presentation_brief["slide_outlines"]:
        evidence_units = evidence_lookup.get(outline["slide_id"], [])
        supporting_points = [unit["value"] for unit in evidence_units] or (outline.get("key_points") or [presentation_brief["objective"]])
        if outline["intent"] in {"risks", "closing"} and presentation_brief.get("known_gaps"):
            supporting_points = [
                *supporting_points,
                *presentation_brief["known_gaps"][:1],
            ]
        if outline["intent"] == "closing" and presentation_brief.get("external_research_questions"):
            supporting_points = [
                *supporting_points,
                f"Next research: {presentation_brief['external_research_questions'][0]}",
            ]
        if outline["intent"] in {"risks", "closing"} and contradiction_summaries:
            supporting_points = [
                *supporting_points,
                f"Conflicted evidence: {contradiction_summaries[0]}",
            ]
        previous_title = story_slides[-1]["title"] if story_slides else None
        transition = _transition_text(previous_title, outline["title"], narrative_pattern)
        claim_posture_summary = _claim_posture_summary(evidence_units)
        validation_note_summary = _validation_note_summary(evidence_units)
        gap_summary = (
            f"Keep these proof gaps visible: {'; '.join(presentation_brief.get('known_gaps', [])[:2])}."
            if presentation_brief.get("known_gaps")
            else ""
        )
        contradiction_summary = (
            f"Call out conflicted evidence explicitly: {'; '.join(contradiction_summaries[:2])}."
            if contradiction_summaries
            else ""
        )
        base_speaker_notes = outline.get("speaker_notes") or f"Connect {outline['title']} back to the decision required and source evidence."
        story_slides.append(
            {
                "slide_id": outline["slide_id"],
                "title": outline["title"],
                "story_role": _story_role_for_intent(outline["intent"]),
                "headline": outline.get("core_message") or outline["title"],
                "core_message": outline.get("core_message") or presentation_brief["narrative_goal"],
                "supporting_points": supporting_points,
                "evidence_refs": [unit["source_ref"] for unit in evidence_units] or outline.get("evidence_refs") or presentation_brief["source_artifact_ids"],
                "speaker_notes": " ".join(
                    part
                    for part in [
                        base_speaker_notes,
                        claim_posture_summary,
                        validation_note_summary,
                        gap_summary,
                        contradiction_summary,
                    ]
                    if part
                ),
                "visual_direction": outline.get("visual_direction") or "Use restrained hierarchy with visible source labels.",
                "confidence_state": evidence_units[0]["confidence_state"] if evidence_units else _confidence_for_intent(outline["intent"]),
                "audience_question_answered": outline["audience_question"],
                "claim": outline["claim"],
                "why_now": presentation_brief["success_outcome"],
                "proof_strategy": "; ".join(
                    [
                        *outline["proof_requirements"],
                        *(
                            [f"Claim posture: {claim_posture_summary}"]
                            if claim_posture_summary
                            else []
                        ),
                        *(
                            [f"Open research: {presentation_brief['external_research_questions'][0]}"]
                            if outline["intent"] in {"risks", "closing"} and presentation_brief.get("external_research_questions")
                            else []
                        ),
                        *(
                            [f"Contradictions: {contradiction_summaries[0]}"]
                            if outline["intent"] in {"risks", "closing"} and contradiction_summaries
                            else []
                        ),
                    ]
                ),
                "transition_from_previous": transition,
                "cta": outline.get("cta") or presentation_brief["decision_required"],
                "appendix_reference": outline.get("appendix_link") or f"appendix_{outline['slide_id']}",
            }
        )

    if contradiction_summaries:
        previous_title = story_slides[-1]["title"] if story_slides else None
        story_slides.append(
            {
                "slide_id": "slide_conflicted_evidence",
                "title": "Conflicted Evidence",
                "story_role": "appendix",
                "headline": "External evidence is directionally useful, but not fully settled.",
                "core_message": "Keep the decision bounded and visible because external sources still disagree on important proof or posture questions.",
                "supporting_points": contradiction_summaries[:3],
                "evidence_refs": presentation_brief["source_artifact_ids"],
                "speaker_notes": " ".join(
                    [
                        "Use this slide when stakeholders need to see where the evidence is still contested.",
                        f"Conflicted evidence: {'; '.join(contradiction_summaries[:2])}.",
                        "Treat this as a trust-preserving appendix or risk bridge, not as a hidden note.",
                    ]
                ),
                "visual_direction": "Evidence appendix with contradiction callouts and explicit source labels.",
                "confidence_state": "low",
                "audience_question_answered": "Where does the external evidence still disagree enough to require bounded decision-making?",
                "claim": "The current recommendation should stay bounded because some external evidence is still contested.",
                "why_now": presentation_brief["success_outcome"],
                "proof_strategy": "; ".join(
                    [
                        "Name the contradiction clearly",
                        "Show the conflicting evidence side by side",
                        "Keep the next review or research action explicit",
                    ]
                ),
                "transition_from_previous": _transition_text(previous_title, "Conflicted Evidence", narrative_pattern),
                "cta": "Resolve the conflicted evidence before broadening the external claim set.",
                "appendix_reference": "appendix_conflicted_evidence",
            }
        )

    return {
        "schema_version": "1.0.0",
        "presentation_story_id": f"presentation_story_{presentation_brief['presentation_brief_id']}",
        "workspace_id": presentation_brief["workspace_id"],
        "presentation_brief_id": presentation_brief["presentation_brief_id"],
        "evidence_pack_id": evidence_pack["evidence_pack_id"],
        "audience": presentation_brief["audience"],
        "presentation_mode": presentation_brief["presentation_mode"],
        "primary_question": presentation_brief["objective"],
        "decision_required": presentation_brief["decision_required"],
        "recommendation": presentation_brief["narrative_goal"],
        "narrative_pattern": narrative_pattern,
        "audience_key_question": presentation_brief["objective"],
        "opening_answer": presentation_brief["narrative_goal"],
        "main_recommendation": presentation_brief["decision_required"],
        "objections_addressed": presentation_brief.get("required_objections", []),
        "appendix_strategy": presentation_brief["appendix_mode"],
        "story_arc": _story_arc_for_brief(presentation_brief),
        "story_transitions": [
            slide["transition_from_previous"] for slide in story_slides
        ],
        "slides": story_slides,
        "appendix_candidates": [
            "Detailed evidence trace for source artifacts.",
            "Fallback PPT composition notes.",
            *(
                ["Conflicted evidence appendix with unresolved external contradictions."]
                if contradiction_summaries
                else []
            ),
        ],
        "generated_at": now_iso(),
    }


def build_render_spec(
    presentation_brief: dict[str, Any],
    presentation_story: dict[str, Any],
    aspect_ratio: str = "16:9",
) -> dict[str, Any]:
    slides = []
    for story_slide in presentation_story["slides"]:
        outline = next(
            (
                outline
                for outline in presentation_brief["slide_outlines"]
                if outline["slide_id"] == story_slide["slide_id"]
            ),
            None,
        )
        if outline is None:
            outline = _synthetic_outline_for_story_slide(presentation_brief, story_slide)
        composition_type = _composition_for_intent(outline["intent"])
        composition_payload = _build_composition_payload(outline, {**story_slide, "composition_type": composition_type})
        chart_spec = _build_chart_spec_for_slide(
            presentation_brief,
            outline,
            story_slide,
            composition_type,
            composition_payload,
        )
        if chart_spec:
            composition_payload["chart_spec"] = chart_spec
        html_target_profile = _html_target_profile(presentation_brief, composition_type)
        ppt_target_profile = _ppt_target_profile(presentation_brief, composition_type, html_target_profile)
        slides.append(
            _apply_render_budgets(
                {
                    "slide_id": story_slide["slide_id"],
                    "title": story_slide["title"],
                    "composition_type": composition_type,
                    "headline": story_slide["headline"],
                    "core_message": story_slide["core_message"],
                    "composition_payload": composition_payload,
                    "layout_variant": _layout_variant(composition_type, presentation_brief["presentation_mode"]),
                    "density_mode": _density_mode(presentation_brief),
                    "visual_tokens": _visual_tokens(presentation_brief, composition_type),
                    "visibility_rules": _visibility_rules(presentation_brief, outline),
                    "annotation_rules": _annotation_rules(outline),
                    "html_render_hints": {
                        "prefer_native_shapes": True,
                        "fallback_layout": "stacked_panels" if composition_type in {"comparison_table", "evidence_grid"} else "standard",
                        "target_profile": html_target_profile,
                        "emphasize_claim": True,
                        "show_evidence_as_cards": composition_type in {"summary_cards", "evidence_grid", "comparison_table"},
                    },
                    "ppt_render_hints": {
                        "prefer_native_shapes": True,
                        "fallback_layout": "stacked_panels" if composition_type in {"comparison_table", "evidence_grid"} else "standard",
                        "target_profile": ppt_target_profile,
                    },
                    "source_refs": story_slide["evidence_refs"],
                    "speaker_notes": story_slide["speaker_notes"],
                    "visual_direction": story_slide["visual_direction"],
                    "confidence_state": story_slide["confidence_state"],
                    "variant_rules": {
                        "live": "Favor the headline and keep most support in notes.",
                        "async": "Keep supporting evidence visible with source labels.",
                        "memo": "Expand the message into short explanatory prose.",
                    },
                },
                presentation_brief,
            )
        )

    return {
        "schema_version": "1.0.0",
        "render_spec_id": f"render_spec_{presentation_brief['presentation_brief_id']}",
        "presentation_story_id": presentation_story["presentation_story_id"],
        "presentation_brief_id": presentation_brief["presentation_brief_id"],
        "workspace_id": presentation_brief["workspace_id"],
        "presentation_mode": presentation_brief["presentation_mode"],
        "theme": theme_preset(presentation_brief["theme_preset"]),
        "aspect_ratio": aspect_ratio,
        "slides": slides,
        "generated_at": now_iso(),
    }


def build_publish_check(
    presentation_brief: dict[str, Any],
    render_spec: dict[str, Any],
    target_format: str = "html",
) -> dict[str, Any]:
    blocking_issues = []
    claim_support_exceptions = []
    density_exceptions = []
    dual_target_required = presentation_brief.get("presentation_format", "both") == "both"
    for slide in render_spec["slides"]:
        html_target_profile = _slide_html_target_profile(slide)
        ppt_target_profile = _slide_ppt_target_profile(slide)
        if not slide["source_refs"]:
            blocking_issues.append(f"{slide['slide_id']} is missing source references.")
        if slide["confidence_state"] == "low" and "risk" not in slide["composition_type"]:
            blocking_issues.append(f"{slide['slide_id']} has low confidence without explicit risk treatment.")
        if not slide["composition_payload"]:
            blocking_issues.append(f"{slide['slide_id']} is missing composition payload.")
        if not _payload_evidence_items(slide):
            claim_support_exceptions.append(f"{slide['slide_id']} has no evidence items attached to its primary claim.")
        if slide["density_mode"] == "dense" and len(_payload_evidence_items(slide)) < 2:
            density_exceptions.append(f"{slide['slide_id']} is marked dense without enough evidence blocks to justify that density.")
        if slide["composition_payload"].get("primary_claim") == slide["title"]:
            claim_support_exceptions.append(f"{slide['slide_id']} primary claim still mirrors the title rather than a sharper claim.")
        if len(slide["headline"]) > _HEADLINE_CHAR_BUDGET.get(slide["composition_type"], 88):
            density_exceptions.append(f"{slide['slide_id']} headline exceeds the layout budget for {slide['composition_type']}.")
        if len(slide["core_message"]) > _CORE_MESSAGE_CHAR_BUDGET.get(slide["composition_type"], 120):
            density_exceptions.append(f"{slide['slide_id']} core message exceeds the layout budget for {slide['composition_type']}.")
        if len(_payload_evidence_items(slide)) > _ITEM_COUNT_BUDGET.get(slide["composition_type"], 4):
            density_exceptions.append(f"{slide['slide_id']} carries too many evidence blocks for a one-screen render.")
        if not _slide_has_visual_payload(slide):
            density_exceptions.append(f"{slide['slide_id']} lacks composition data needed for a meaningful visual treatment.")
        if dual_target_required and html_target_profile != "dual_target":
            blocking_issues.append(
                f"{slide['slide_id']} is marked {html_target_profile} for HTML and cannot ship in a both-format deck."
            )
        if dual_target_required and ppt_target_profile != "dual_target":
            blocking_issues.append(
                f"{slide['slide_id']} is marked {ppt_target_profile} for PPT and cannot ship in a both-format deck."
            )
        if target_format == "pptx" and html_target_profile == "html_rich" and ppt_target_profile != "ppt_safe":
            blocking_issues.append(
                f"{slide['slide_id']} uses an HTML-rich design without an explicit PPT-safe fallback."
            )

    fidelity_exceptions = []
    if target_format == "pptx":
        ppt_safe_slides = [
            slide["slide_id"] for slide in render_spec["slides"] if _slide_ppt_target_profile(slide) == "ppt_safe"
        ]
        if ppt_safe_slides:
            fidelity_exceptions.append(
                "PPT-safe fallback rendering is active for: " + ", ".join(ppt_safe_slides) + "."
            )

    narrative_quality_score, narrative_notes = _narrative_quality(render_spec, presentation_brief)
    evidence_quality_score, evidence_notes = _evidence_quality(render_spec)
    audience_fit_score, audience_notes = _audience_fit_quality(render_spec, presentation_brief)
    visual_clarity_score, visual_notes = _visual_clarity_quality(render_spec)
    missing_objections = [] if presentation_brief.get("required_objections") else ["No required objections were specified in the brief."]

    if narrative_quality_score < 3.2:
        blocking_issues.append("Narrative quality is below the minimum publish bar.")
    if evidence_quality_score < 3.2:
        blocking_issues.append("Evidence quality is below the minimum publish bar.")
    if audience_fit_score < 3.0:
        blocking_issues.append("Audience fit is below the minimum publish bar.")
    if visual_clarity_score < 3.0:
        blocking_issues.append("Visual clarity is below the minimum publish bar.")

    html_profiles = {_slide_html_target_profile(slide) for slide in render_spec["slides"]}
    ppt_profiles = {_slide_ppt_target_profile(slide) for slide in render_spec["slides"]}
    html_fidelity_status = "aligned" if html_profiles == {"dual_target"} else "partial"
    if dual_target_required and "html_rich" in html_profiles:
        html_fidelity_status = "at_risk"
    ppt_fidelity_status = "aligned" if ppt_profiles == {"dual_target"} else "planned_fallbacks"
    if dual_target_required and "ppt_safe" in ppt_profiles:
        ppt_fidelity_status = "at_risk"

    return {
        "schema_version": "1.0.0",
        "publish_check_id": f"publish_check_{presentation_brief['presentation_brief_id']}_{target_format}",
        "workspace_id": presentation_brief["workspace_id"],
        "presentation_brief_id": presentation_brief["presentation_brief_id"],
        "render_spec_id": render_spec["render_spec_id"],
        "target_format": target_format,
        "publish_ready": not blocking_issues,
        "blocking_issues": blocking_issues,
        "evidence_coverage": "full" if all(slide["source_refs"] for slide in render_spec["slides"]) else "partial",
        "risk_visibility": "explicit"
        if any(slide["composition_type"] == "risk_matrix" for slide in render_spec["slides"])
        else "partial",
        "audience_safety_status": "safe",
        "narrative_quality_score": narrative_quality_score,
        "evidence_quality_score": evidence_quality_score,
        "audience_fit_score": audience_fit_score,
        "visual_clarity_score": visual_clarity_score,
        "html_fidelity_status": html_fidelity_status,
        "ppt_fidelity_status": ppt_fidelity_status,
        "redaction_status": "safe" if not presentation_brief["customer_safe"] or presentation_brief["redaction_policy"] != "none" else "mixed",
        "claim_support_exceptions": claim_support_exceptions + narrative_notes + evidence_notes,
        "density_exceptions": density_exceptions + audience_notes + visual_notes,
        "missing_objections": missing_objections,
        "fidelity_exceptions": fidelity_exceptions,
        "generated_at": now_iso(),
    }


def _render_supporting_list(items: list[dict[str, Any]]) -> str:
    cards = []
    for item in items:
        evidence_kind = escape(item.get("evidence_type", "fact"))
        cards.append(
            f"<article class='evidence-card evidence-{evidence_kind}'>"
            f"<div class='evidence-label-row'><div class='evidence-label'>{escape(item['label'])}</div><div class='evidence-kind'>{evidence_kind.replace('_', ' ')}</div></div>"
            f"<div class='evidence-text'>{escape(item['value'])}</div>"
            f"<div class='evidence-annotation'>{escape(item['annotation'])}</div>"
            f"<div class='source-tag'>{escape(item['source_ref'])}</div>"
            "</article>"
        )
    return "".join(cards)


def _render_decision_options(options: list[dict[str, Any]]) -> str:
    if not options:
        return ""
    cards = []
    for option in options:
        tradeoffs = "".join(f"<li>{escape(item)}</li>" for item in option["tradeoffs"])
        cards.append(
            "<article class='decision-option'>"
            f"<div class='decision-option-label'>{escape(option['label'])}</div>"
            f"<div class='decision-option-summary'>{escape(option['summary'])}</div>"
            f"<ul class='decision-option-tradeoffs'>{tradeoffs}</ul>"
            "</article>"
        )
    return f"<div class='decision-options'>{''.join(cards)}</div>"


def _render_risk_items(risk_items: list[dict[str, Any]]) -> str:
    if not risk_items:
        return ""
    cards = []
    for item in risk_items:
        cards.append(
            "<article class='risk-item'>"
            f"<div class='risk-title'>{escape(item['label'])}</div>"
            f"<div class='risk-meta'>{escape(item['likelihood'])} likelihood / {escape(item['impact'])} impact</div>"
            f"<div class='risk-owner'>Owner: {escape(item['owner'])}</div>"
            f"<div class='risk-mitigation'>{escape(item['mitigation'])}</div>"
            "</article>"
        )
    return f"<div class='risk-matrix'>{''.join(cards)}</div>"


def _render_timeline_events(events: list[dict[str, Any]]) -> str:
    if not events:
        return ""
    items = []
    for event in events:
        items.append(
            "<article class='timeline-event'>"
            f"<div class='timeline-timing'>{escape(event['timing'])}</div>"
            f"<div class='timeline-label'>{escape(event['label'])}</div>"
            f"<div class='timeline-dependency'>{escape(event['dependency'])}</div>"
            "</article>"
        )
    return f"<div class='timeline-track'>{''.join(items)}</div>"


def _render_comparison_rows(rows: list[dict[str, Any]]) -> str:
    if not rows:
        return ""
    header = (
        "<div class='comparison-row comparison-header'>"
        "<span>Dimension</span><span>Current</span><span>Target</span><span>Highlight</span>"
        "</div>"
    )
    body = []
    for row in rows:
        body.append(
            "<div class='comparison-row'>"
            f"<span>{escape(row['dimension'])}</span>"
            f"<span>{escape(row['current_state'])}</span>"
            f"<span>{escape(row['target_state'])}</span>"
            f"<span>{escape(row['highlight'])}</span>"
            "</div>"
        )
    return f"<div class='comparison-table'>{header}{''.join(body)}</div>"


def _render_metric_strip(rows: list[dict[str, Any]]) -> str:
    if not rows:
        return ""
    cards = []
    for row in rows:
        cards.append(
            "<article class='metric-card'>"
            f"<div class='metric-label'>{escape(row['dimension'])}</div>"
            f"<div class='metric-value'>{escape(row['current_state'])}</div>"
            f"<div class='metric-note'>{escape(row['highlight'])}</div>"
            "</article>"
        )
    return f"<div class='metric-strip'>{''.join(cards)}</div>"


def _render_appendix_trace(rows: list[dict[str, Any]]) -> str:
    if not rows:
        return ""
    items = []
    for row in rows:
        items.append(
            "<article class='appendix-trace-item'>"
            f"<div class='appendix-trace-source'>{escape(row['current_state'])}</div>"
            f"<div class='appendix-trace-value'>{escape(row['target_state'])}</div>"
            f"<div class='appendix-trace-note'>{escape(row['highlight'])}</div>"
            "</article>"
        )
    return f"<div class='appendix-trace'>{''.join(items)}</div>"


def _render_hero_signal_rail(items: list[dict[str, Any]]) -> str:
    if not items:
        return ""
    bars = []
    for index, item in enumerate(items[:3], start=1):
        width = 44 + index * 14
        bars.append(
            "<article class='hero-signal'>"
            f"<div class='hero-signal-label'>{escape(item['label'])}</div>"
            f"<div class='hero-signal-bar'><span style='width:{width}%;'></span></div>"
            f"<div class='hero-signal-text'>{escape(item['value'])}</div>"
            "</article>"
        )
    return f"<div class='hero-signal-rail'>{''.join(bars)}</div>"


def _risk_cell_position(likelihood: str, impact: str) -> tuple[int, int]:
    likelihood_pos = {"low": 0, "medium": 1, "high": 2}
    impact_pos = {"low": 2, "medium": 1, "high": 0}
    return likelihood_pos.get(likelihood, 1), impact_pos.get(impact, 1)


def _render_risk_heatmap(risk_items: list[dict[str, Any]]) -> str:
    if not risk_items:
        return ""
    dots = []
    labels = []
    for index, item in enumerate(risk_items[:4], start=1):
        col, row = _risk_cell_position(item["likelihood"], item["impact"])
        dots.append(
            f"<span class='risk-dot risk-dot-{index}' style='grid-column:{col + 1}; grid-row:{row + 1};'>{index}</span>"
        )
        labels.append(
            "<div class='risk-legend-item'>"
            f"<span class='risk-legend-index'>{index}</span>"
            f"<span>{escape(item['label'])}</span>"
            "</div>"
        )
    return (
        "<div class='risk-heatmap-panel'>"
        "<div class='risk-heatmap'>"
        "<span class='risk-axis risk-axis-y'>Impact</span>"
        "<span class='risk-axis risk-axis-x'>Likelihood</span>"
        f"{''.join(dots)}"
        "</div>"
        f"<div class='risk-legend'>{''.join(labels)}</div>"
        "</div>"
    )


def _render_comparison_visual(rows: list[dict[str, Any]]) -> str:
    if not rows:
        return ""
    cards = []
    for row in rows[:3]:
        cards.append(
            "<article class='comparison-visual-card'>"
            f"<div class='comparison-visual-dimension'>{escape(row['dimension'])}</div>"
            "<div class='comparison-visual-lanes'>"
            f"<div class='comparison-lane comparison-lane-current'><span>Now</span><strong>{escape(row['current_state'])}</strong></div>"
            f"<div class='comparison-lane comparison-lane-target'><span>Next</span><strong>{escape(row['target_state'])}</strong></div>"
            "</div>"
            f"<div class='comparison-visual-highlight'>{escape(row['highlight'])}</div>"
            "</article>"
        )
    return f"<div class='comparison-visual-grid'>{''.join(cards)}</div>"


def _render_chart_container(slide: dict[str, Any]) -> str:
    chart_spec = slide["composition_payload"].get("chart_spec")
    if not chart_spec:
        return ""
    heuristic_note = (
        "<div class='chart-footnote'>Heuristic signal index from qualitative evidence.</div>"
        if chart_spec.get("uses_heuristics")
        else ""
    )
    return (
        "<div class='chart-panel'>"
        f"<div class='chart-header'><div class='chart-title'>{escape(chart_spec['title'])}</div>"
        f"<div class='chart-subtitle'>{escape(chart_spec.get('subtitle', ''))}</div></div>"
        f"<div class='echart-canvas' id='{escape(chart_spec['chart_id'])}'></div>"
        "<div class='chart-fallback'>Interactive chart loads when ECharts is available.</div>"
        f"{heuristic_note}"
        "</div>"
    )


def _render_map_item(item: dict[str, Any]) -> str:
    description = f"<div class='map-item-description'>{escape(item.get('description', ''))}</div>" if item.get("description") else ""
    return (
        "<div class='map-item-card'>"
        f"<div class='map-item-label'>{escape(item['label'])}</div>"
        f"{description}"
        "</div>"
    )


def _render_map_canvas(slide: dict[str, Any]) -> str:
    map_payload = slide["composition_payload"].get("map_payload") or {}
    map_type = map_payload.get("map_type", slide["composition_type"])

    if map_type in _MATRIX_MAP_COMPOSITIONS:
        axis_labels = map_payload.get("axes", [])
        items = map_payload.get("items", [])
        cells = []
        for index in range(4):
            item = items[index] if index < len(items) else {"label": "Open quadrant", "description": "No mapped item yet."}
            extra_class = " matrix-cell-empty" if index >= len(items) else ""
            cells.append(
                f"<div class='matrix-cell{extra_class}'>"
                f"<div class='matrix-cell-index'>{index + 1:02d}</div>"
                f"{_render_map_item(item)}"
                "</div>"
            )
        certainty = "".join(
            f"<span class='summary-chip'>{escape(note)}</span>" for note in map_payload.get("certainty_notes", [])[:3]
        )
        return (
            "<div class='map-layout'>"
            "<div class='matrix-map'>"
            f"<div class='matrix-axis matrix-axis-y'>{escape(axis_labels[0] if len(axis_labels) > 0 else 'Axis 1')}</div>"
            f"<div class='matrix-axis matrix-axis-x'>{escape(axis_labels[1] if len(axis_labels) > 1 else 'Axis 2')}</div>"
            f"{''.join(cells)}"
            "</div>"
            f"<div class='summary-chip-row'>{certainty}</div>"
            "</div>"
        )

    stage_markup = ""
    if map_payload.get("stages"):
        stage_markup = "".join(
            f"<div class='map-stage'>{escape(stage)}</div>" for stage in map_payload["stages"][:6]
        )
        stage_markup = f"<div class='map-stage-row'>{stage_markup}</div>"

    lane_markup = ""
    if map_payload.get("lanes"):
        lane_markup = "".join(
            f"<div class='map-lane'>{escape(lane)}</div>" for lane in map_payload["lanes"][:6]
        )
        lane_markup = f"<div class='map-lane-row'>{lane_markup}</div>"

    node_items = map_payload.get("items") or map_payload.get("nodes") or []
    node_markup = "".join(_render_map_item(item) for item in node_items[:6])
    axes_markup = ""
    if map_payload.get("axes"):
        axes_markup = "".join(
            f"<span class='source-tag map-axis-tag'>{escape(axis)}</span>" for axis in map_payload["axes"][:4]
        )
        axes_markup = f"<div class='map-axis-row'>{axes_markup}</div>"
    certainty = "".join(
        f"<span class='summary-chip'>{escape(note)}</span>" for note in map_payload.get("certainty_notes", [])[:3]
    )

    return (
        "<div class='map-layout'>"
        f"{stage_markup}"
        f"{lane_markup}"
        f"{axes_markup}"
        f"<div class='map-grid map-type-{escape(map_type)}'>{node_markup}</div>"
        f"<div class='summary-chip-row'>{certainty}</div>"
        "</div>"
    )


def _render_slide_body(slide: dict[str, Any]) -> str:
    payload = slide["composition_payload"]
    evidence_html = _render_supporting_list(_payload_evidence_items(slide))

    if slide["composition_type"] == "hero_statement":
        recommendation = payload.get("recommendation")
        decision_ask = payload.get("decision_ask")
        strip = ""
        if recommendation or decision_ask:
            strip = (
                "<div class='hero-strip'>"
                f"<div><span class='strip-label'>Recommendation</span><strong>{escape(recommendation or '')}</strong></div>"
                f"<div><span class='strip-label'>Ask</span><strong>{escape(decision_ask or '')}</strong></div>"
                "</div>"
            )
        spotlight = (
            "<div class='hero-spotlight'>"
            f"<div class='hero-claim'>{escape(payload.get('primary_claim', slide['core_message']))}</div>"
            f"<div class='hero-summary'>{escape(payload.get('summary', slide['core_message']))}</div>"
            "</div>"
        )
        visual = _render_hero_signal_rail(payload.get("evidence_items", []))
        return f"<div class='hero-layout'>{spotlight}{visual}{strip}<div class='evidence-grid compact-grid'>{evidence_html}</div></div>"

    if slide["composition_type"] == "decision_frame":
        return (
            "<div class='decision-layout'>"
            f"<div class='decision-summary'><div class='decision-ask'>{escape(payload.get('decision_ask', ''))}</div>{_render_decision_options(payload.get('options', []))}</div>"
            f"<div class='decision-proof'>{evidence_html}</div>"
            "</div>"
        )

    if slide["composition_type"] == "risk_matrix":
        return (
            "<div class='risk-layout'>"
            f"{_render_risk_heatmap(payload.get('risk_items', []))}"
            f"{_render_risk_items(payload.get('risk_items', []))}"
            f"<div class='evidence-grid'>{evidence_html}</div>"
            "</div>"
        )

    if slide["composition_type"] == "timeline_with_dependencies":
        return (
            "<div class='timeline-layout'>"
            f"{_render_timeline_events(payload.get('timeline_events', []))}"
            f"<div class='evidence-grid compact-grid'>{evidence_html}</div>"
            "</div>"
        )

    if slide["composition_type"] == "comparison_table":
        return (
            "<div class='comparison-layout'>"
            f"{_render_chart_container(slide)}"
            f"{_render_comparison_visual(payload.get('comparison_rows', []))}"
            f"{_render_comparison_rows(payload.get('comparison_rows', []))}"
            f"<div class='evidence-grid compact-grid'>{evidence_html}</div>"
            "</div>"
        )

    if slide["composition_type"] == "metric_strip":
        return (
            "<div class='metric-layout'>"
            f"{_render_chart_container(slide)}"
            f"{_render_metric_strip(payload.get('comparison_rows', []))}"
            f"<div class='evidence-grid compact-grid'>{evidence_html}</div>"
            "</div>"
        )

    if slide["composition_type"] == "appendix_evidence":
        return (
            "<div class='appendix-layout'>"
            f"{_render_appendix_trace(payload.get('comparison_rows', []))}"
            f"<div class='evidence-grid compact-grid'>{evidence_html}</div>"
            "</div>"
        )

    if slide["composition_type"] == "evidence_grid":
        return f"<div class='evidence-grid'>{evidence_html}</div>"

    if slide["composition_type"] == "summary_cards":
        summary_tags = "".join(
            f"<span class='summary-chip'>{escape(item['label'])}</span>" for item in _payload_evidence_items(slide)[:4]
        )
        return f"<div class='summary-layout'><div class='summary-chip-row'>{summary_tags}</div><div class='evidence-grid'>{evidence_html}</div></div>"

    if slide["composition_type"] in _MAP_COMPOSITIONS:
        return (
            "<div class='map-body-layout'>"
            f"{_render_map_canvas(slide)}"
            f"<div class='evidence-grid compact-grid'>{evidence_html}</div>"
            "</div>"
        )

    return f"<div class='evidence-grid'>{evidence_html}</div>"


def render_render_spec_html(render_spec: dict[str, Any]) -> str:
    theme = render_spec["theme"]
    slide_sections = []
    chart_payloads = []
    for index, slide in enumerate(render_spec["slides"], start=1):
        chart_spec = slide["composition_payload"].get("chart_spec")
        if chart_spec:
            chart_payloads.append(chart_spec)
        slide_sections.append(
            f"""
            <section class="slide composition-{slide['composition_type']} density-{slide['density_mode']} layout-{slide['layout_variant']} profile-{slide['html_render_hints'].get('target_profile', 'dual_target')}">
              <div class="slide-stage">
                <div class="slide-frame">
                  <div class="slide-content">
                    <div class="slide-number">{index:02d}</div>
                    <div class="slide-topline">
                      <span class="eyebrow">{escape(render_spec['presentation_mode'])} / {escape(slide['title'])}</span>
                      <div class="topline-badges">
                        <span class="badge badge-mode">{escape(slide['layout_variant']).replace('_', ' ')}</span>
                        <span class="confidence confidence-{slide['confidence_state']}">{escape(slide['confidence_state'])} confidence</span>
                      </div>
                    </div>
                    <div class="headline-block">
                      <h2>{escape(slide['headline'])}</h2>
                      <p class="core-message">{escape(slide['core_message'])}</p>
                      <p class="primary-claim">{escape(slide['composition_payload'].get('primary_claim', slide['core_message']))}</p>
                    </div>
                    <div class="slide-body">{_render_slide_body(slide)}</div>
                    <div class="slide-footer">
                      <div class="visual-direction">{escape(slide['visual_direction'])}</div>
                      <div class="source-rail">Sources: {escape(", ".join(slide['source_refs']))}</div>
                    </div>
                  </div>
                </div>
              </div>
              <details class="slide-notes">
                <summary>Speaker Notes</summary>
                <div class="speaker-notes">{escape(slide['speaker_notes'])}</div>
              </details>
            </section>
            """
        )

    return f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{escape(render_spec['render_spec_id'])}</title>
  <script defer src="https://cdn.jsdelivr.net/npm/echarts@5/dist/echarts.min.js"></script>
  <style>
    :root {{
      --font-family: {theme['font_family']};
      --background: {theme['background']};
      --accent: {theme['accent']};
      --ink: #111111;
      --paper: rgba(255,255,255,0.78);
      --line: rgba(17,17,17,0.10);
      --panel-shadow: 0 24px 80px rgba(0,0,0,0.12);
      --grain:
        linear-gradient(120deg, rgba(255,255,255,0.22), rgba(255,255,255,0.02)),
        radial-gradient(circle at 20% 20%, rgba(255,255,255,0.18), transparent 36%),
        radial-gradient(circle at 80% 0%, rgba(255,255,255,0.12), transparent 28%);
    }}
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      font-family: var(--font-family);
      background: var(--background);
      background-attachment: fixed;
      color: var(--ink);
      position: relative;
    }}
    body::before {{
      content: "";
      position: fixed;
      inset: 0;
      background-image:
        linear-gradient(rgba(255,255,255,0.06) 1px, transparent 1px),
        linear-gradient(90deg, rgba(255,255,255,0.06) 1px, transparent 1px);
      background-size: 36px 36px;
      opacity: 0.24;
      pointer-events: none;
    }}
    .deck {{
      display: grid;
      gap: 36px;
      padding: 24px 24px 36px;
      max-width: none;
      margin: 0 auto;
    }}
    .slide {{
      --slide-stage-width: min(1320px, calc(100vw - 56px));
      display: grid;
      gap: 12px;
      position: relative;
      justify-items: center;
    }}
    .slide-stage {{
      width: var(--slide-stage-width);
      max-width: 100%;
    }}
    .slide-frame {{
      width: 100%;
      aspect-ratio: 16 / 9;
      min-height: 0;
      border-radius: 28px;
      background:
        var(--grain),
        linear-gradient(180deg, rgba(255,255,255,0.88), rgba(255,255,255,0.76));
      border: 1px solid rgba(255,255,255,0.45);
      box-shadow: var(--panel-shadow);
      backdrop-filter: blur(10px);
      position: relative;
      isolation: isolate;
      overflow: hidden;
    }}
    .slide-frame::after {{
      content: "";
      position: absolute;
      inset: 0;
      border-radius: inherit;
      background:
        radial-gradient(circle at 100% 0%, rgba(30,90,168,0.10), transparent 32%),
        radial-gradient(circle at 0% 100%, rgba(255,111,60,0.08), transparent 28%);
      pointer-events: none;
      z-index: 0;
    }}
    .slide-content {{
      --fit-scale: 1;
      position: relative;
      z-index: 1;
      width: calc(100% / var(--fit-scale));
      min-height: calc(100% / var(--fit-scale));
      padding: 24px 28px 20px;
      display: grid;
      grid-template-rows: auto auto minmax(0, 1fr) auto;
      gap: 14px;
      align-content: start;
      transform: scale(var(--fit-scale));
      transform-origin: top left;
    }}
    .slide-topline,
    .headline-block,
    .slide-body,
    .slide-footer,
    .slide-number {{
      position: relative;
      z-index: 1;
    }}
    .slide-number {{
      position: absolute;
      top: 18px;
      right: 22px;
      font-size: 11px;
      letter-spacing: 0.16em;
      text-transform: uppercase;
      color: rgba(17,17,17,0.45);
    }}
    .slide-topline {{
      display: flex;
      justify-content: space-between;
      gap: 16px;
      align-items: center;
      flex-wrap: wrap;
    }}
    .headline-block {{
      display: grid;
      gap: 10px;
      max-width: min(66%, 760px);
    }}
    .eyebrow {{
      text-transform: uppercase;
      letter-spacing: 0.18em;
      font-size: 12px;
      color: var(--accent);
    }}
    .topline-badges {{
      display: flex;
      flex-wrap: wrap;
      gap: 8px;
      align-items: center;
    }}
    .badge {{
      padding: 8px 12px;
      border-radius: 999px;
      font-size: 11px;
      letter-spacing: 0.06em;
      text-transform: uppercase;
      background: rgba(17,17,17,0.05);
    }}
    .badge-mode {{
      color: rgba(17,17,17,0.72);
    }}
    .confidence {{
      padding: 8px 12px;
      border-radius: 999px;
      font-size: 12px;
      background: rgba(17,17,17,0.06);
    }}
    .confidence-medium {{ background: rgba(218, 152, 35, 0.16); }}
    .confidence-low {{ background: rgba(176, 43, 43, 0.14); }}
    h2 {{
      margin: 0;
      font-size: 32px;
      line-height: 1.04;
      max-width: 19ch;
      letter-spacing: -0.04em;
    }}
    .core-message {{
      margin: 0;
      font-size: 17px;
      line-height: 1.32;
      max-width: 46rem;
    }}
    .primary-claim {{
      margin: 0;
      font-size: 12px;
      text-transform: uppercase;
      letter-spacing: 0.08em;
      color: var(--accent);
    }}
    .evidence-grid {{
      display: grid;
      grid-template-columns: repeat(2, minmax(0, 1fr));
      gap: 12px;
    }}
    .compact-grid {{
      grid-template-columns: repeat(auto-fit, minmax(220px, 1fr));
    }}
    .evidence-card {{
      border: 1px solid var(--line);
      border-radius: 18px;
      padding: 11px 13px;
      background: rgba(255,255,255,0.56);
      display: grid;
      gap: 8px;
      align-content: start;
      min-height: 0;
    }}
    .evidence-card::before {{
      content: "";
      width: 44px;
      height: 4px;
      border-radius: 999px;
      background: var(--accent);
      opacity: 0.6;
    }}
    .evidence-label-row {{
      display: flex;
      justify-content: space-between;
      gap: 10px;
      align-items: start;
    }}
    .evidence-label {{
      font-size: 11px;
      text-transform: uppercase;
      letter-spacing: 0.12em;
      color: var(--accent);
    }}
    .evidence-kind {{
      font-size: 10px;
      text-transform: uppercase;
      letter-spacing: 0.08em;
      color: rgba(17,17,17,0.52);
      white-space: nowrap;
    }}
    .evidence-text {{
      font-size: 13px;
      line-height: 1.3;
    }}
    .evidence-annotation {{
      font-size: 11px;
      color: rgba(17,17,17,0.65);
    }}
    .source-tag {{
      display: inline-flex;
      width: fit-content;
      padding: 6px 10px;
      border-radius: 999px;
      background: rgba(30,90,168,0.10);
      color: var(--accent);
      font-size: 12px;
    }}
    .slide-body {{
      display: grid;
      gap: 10px;
      align-content: start;
      min-height: 0;
    }}
    .hero-spotlight {{
      display: grid;
      gap: 8px;
      padding: 14px 16px;
      border-radius: 18px;
      background:
        radial-gradient(circle at 0% 0%, rgba(30,90,168,0.14), transparent 45%),
        linear-gradient(135deg, rgba(255,255,255,0.48), rgba(30,90,168,0.06));
      border: 1px solid rgba(30,90,168,0.12);
    }}
    .hero-claim {{
      font-size: 18px;
      font-weight: 700;
      line-height: 1.15;
    }}
    .hero-summary {{
      font-size: 13px;
      line-height: 1.35;
      color: rgba(17,17,17,0.7);
    }}
    .hero-layout {{
      display: grid;
      grid-template-columns: 1.15fr 0.85fr;
      gap: 10px;
      align-content: start;
    }}
    .hero-layout .hero-strip,
    .hero-layout .evidence-grid {{
      grid-column: 1 / -1;
    }}
    .hero-signal-rail {{
      display: grid;
      gap: 10px;
      padding: 14px;
      border-radius: 18px;
      background: rgba(255,255,255,0.5);
      border: 1px solid rgba(30,90,168,0.12);
      align-content: start;
    }}
    .hero-signal {{
      display: grid;
      gap: 6px;
    }}
    .hero-signal-label {{
      font-size: 10px;
      text-transform: uppercase;
      letter-spacing: 0.1em;
      color: var(--accent);
    }}
    .hero-signal-bar {{
      height: 8px;
      border-radius: 999px;
      background: rgba(17,17,17,0.08);
      overflow: hidden;
    }}
    .hero-signal-bar span {{
      display: block;
      height: 100%;
      border-radius: 999px;
      background: linear-gradient(90deg, rgba(30,90,168,0.55), rgba(30,90,168,0.95));
    }}
    .hero-signal-text {{
      font-size: 12px;
      line-height: 1.25;
      color: rgba(17,17,17,0.72);
    }}
    .hero-strip {{
      display: grid;
      grid-template-columns: repeat(2, minmax(0, 1fr));
      gap: 12px;
      padding: 14px 16px;
      border-radius: 18px;
      background: linear-gradient(135deg, rgba(30,90,168,0.12), rgba(255,255,255,0.38));
      border: 1px solid rgba(30,90,168,0.12);
    }}
    .strip-label {{
      display: block;
      font-size: 11px;
      text-transform: uppercase;
      letter-spacing: 0.12em;
      color: var(--accent);
      margin-bottom: 6px;
    }}
    .summary-layout {{
      display: grid;
      gap: 12px;
    }}
    .summary-chip-row {{
      display: flex;
      flex-wrap: wrap;
      gap: 8px;
    }}
    .summary-chip {{
      padding: 6px 10px;
      border-radius: 999px;
      background: rgba(30,90,168,0.09);
      color: var(--accent);
      font-size: 11px;
      text-transform: uppercase;
      letter-spacing: 0.08em;
    }}
    .decision-layout,
    .risk-layout,
    .timeline-layout,
    .comparison-layout,
    .metric-layout,
    .appendix-layout {{
      display: grid;
      gap: 12px;
    }}
    .map-body-layout {{
      display: grid;
      grid-template-columns: 1.15fr 0.85fr;
      gap: 16px;
      align-items: start;
    }}
    .map-layout {{
      display: grid;
      gap: 12px;
      padding: 16px;
      border-radius: 20px;
      background: rgba(255,255,255,0.54);
      border: 1px solid rgba(30,90,168,0.12);
    }}
    .map-stage-row,
    .map-lane-row {{
      display: grid;
      grid-template-columns: repeat(auto-fit, minmax(110px, 1fr));
      gap: 10px;
    }}
    .map-stage,
    .map-lane {{
      padding: 10px 12px;
      border-radius: 14px;
      font-size: 11px;
      text-transform: uppercase;
      letter-spacing: 0.1em;
      color: var(--accent);
      background: rgba(30,90,168,0.08);
      border: 1px solid rgba(30,90,168,0.1);
    }}
    .map-lane {{
      background: rgba(17,17,17,0.05);
      color: rgba(17,17,17,0.7);
      border-color: rgba(17,17,17,0.08);
    }}
    .map-axis-row {{
      display: flex;
      flex-wrap: wrap;
      gap: 8px;
    }}
    .map-axis-tag {{
      background: rgba(17,17,17,0.06);
      color: rgba(17,17,17,0.72);
    }}
    .map-grid {{
      display: grid;
      grid-template-columns: repeat(auto-fit, minmax(180px, 1fr));
      gap: 12px;
    }}
    .map-item-card {{
      min-height: 110px;
      padding: 14px 16px;
      border-radius: 18px;
      background:
        linear-gradient(180deg, rgba(255,255,255,0.82), rgba(255,255,255,0.56)),
        radial-gradient(circle at 0% 0%, rgba(30,90,168,0.08), transparent 42%);
      border: 1px solid var(--line);
      display: grid;
      gap: 8px;
      align-content: start;
    }}
    .map-item-label {{
      font-size: 15px;
      font-weight: 700;
      line-height: 1.2;
    }}
    .map-item-description {{
      font-size: 12px;
      line-height: 1.35;
      color: rgba(17,17,17,0.7);
    }}
    .matrix-map {{
      position: relative;
      display: grid;
      grid-template-columns: repeat(2, minmax(0, 1fr));
      gap: 12px;
      padding: 28px 18px 20px 28px;
      border-radius: 20px;
      background:
        linear-gradient(180deg, rgba(255,255,255,0.82), rgba(255,255,255,0.56)),
        linear-gradient(90deg, rgba(30,90,168,0.04), rgba(17,17,17,0.03));
      border: 1px solid rgba(17,17,17,0.08);
    }}
    .matrix-map::before,
    .matrix-map::after {{
      content: "";
      position: absolute;
      background: rgba(17,17,17,0.08);
    }}
    .matrix-map::before {{
      top: 18px;
      bottom: 18px;
      left: 50%;
      width: 1px;
    }}
    .matrix-map::after {{
      left: 18px;
      right: 18px;
      top: 50%;
      height: 1px;
    }}
    .matrix-cell {{
      min-height: 140px;
      padding: 16px;
      border-radius: 16px;
      background: rgba(255,255,255,0.62);
      border: 1px solid rgba(17,17,17,0.08);
      display: grid;
      gap: 8px;
      align-content: start;
      position: relative;
      z-index: 1;
    }}
    .matrix-cell-empty {{
      opacity: 0.72;
    }}
    .matrix-cell-index {{
      width: 28px;
      height: 28px;
      border-radius: 999px;
      background: rgba(30,90,168,0.10);
      color: var(--accent);
      display: inline-flex;
      align-items: center;
      justify-content: center;
      font-size: 11px;
      font-weight: 700;
    }}
    .matrix-axis {{
      position: absolute;
      font-size: 10px;
      text-transform: uppercase;
      letter-spacing: 0.12em;
      color: rgba(17,17,17,0.55);
      z-index: 1;
    }}
    .matrix-axis-y {{
      left: -4px;
      top: 50%;
      transform: rotate(-90deg) translateY(-50%);
      transform-origin: left top;
    }}
    .matrix-axis-x {{
      right: 18px;
      bottom: -2px;
    }}
    .decision-layout {{
      grid-template-columns: 1.1fr 0.9fr;
    }}
    .decision-summary,
    .decision-proof {{
      display: grid;
      gap: 14px;
    }}
    .decision-ask {{
      padding: 14px 16px;
      border-radius: 18px;
      background: linear-gradient(135deg, rgba(30,90,168,0.14), rgba(30,90,168,0.05));
      font-weight: 700;
    }}
    .decision-options {{
      display: grid;
      gap: 12px;
    }}
    .decision-option {{
      padding: 14px 16px;
      border-radius: 18px;
      border: 1px solid var(--line);
      background: rgba(255,255,255,0.52);
      display: grid;
      gap: 8px;
    }}
    .decision-option-label {{
      text-transform: uppercase;
      letter-spacing: 0.12em;
      font-size: 11px;
      color: var(--accent);
    }}
    .decision-option-summary {{
      font-size: 15px;
      font-weight: 700;
    }}
    .decision-option-tradeoffs {{
      margin: 0;
      padding-left: 18px;
      display: grid;
      gap: 6px;
    }}
    .risk-matrix {{
      display: grid;
      grid-template-columns: repeat(auto-fit, minmax(220px, 1fr));
      gap: 10px;
    }}
    .risk-heatmap-panel {{
      display: grid;
      grid-template-columns: 220px 1fr;
      gap: 12px;
      align-items: start;
    }}
    .risk-heatmap {{
      position: relative;
      min-height: 210px;
      display: grid;
      grid-template-columns: repeat(3, 1fr);
      grid-template-rows: repeat(3, 1fr);
      gap: 8px;
      padding: 18px 18px 28px 28px;
      border-radius: 18px;
      background:
        linear-gradient(180deg, rgba(176,43,43,0.12), rgba(218,152,35,0.08) 55%, rgba(47,125,83,0.08));
      border: 1px solid rgba(17,17,17,0.08);
    }}
    .risk-heatmap::before {{
      content: "";
      position: absolute;
      inset: 18px 18px 28px 28px;
      display: block;
      border-radius: 12px;
      background:
        linear-gradient(rgba(255,255,255,0.45), rgba(255,255,255,0.45)),
        linear-gradient(90deg, transparent 32.8%, rgba(17,17,17,0.1) 33%, rgba(17,17,17,0.1) 34%, transparent 34.2%, transparent 65.8%, rgba(17,17,17,0.1) 66%, rgba(17,17,17,0.1) 67%, transparent 67.2%),
        linear-gradient(transparent 32.8%, rgba(17,17,17,0.1) 33%, rgba(17,17,17,0.1) 34%, transparent 34.2%, transparent 65.8%, rgba(17,17,17,0.1) 66%, rgba(17,17,17,0.1) 67%, transparent 67.2%);
      opacity: 0.9;
    }}
    .risk-dot {{
      align-self: center;
      justify-self: center;
      width: 34px;
      height: 34px;
      border-radius: 50%;
      background: var(--accent);
      color: white;
      font-size: 14px;
      font-weight: 700;
      display: inline-flex;
      align-items: center;
      justify-content: center;
      z-index: 1;
      box-shadow: 0 8px 18px rgba(30,90,168,0.25);
    }}
    .risk-axis {{
      position: absolute;
      font-size: 10px;
      text-transform: uppercase;
      letter-spacing: 0.1em;
      color: rgba(17,17,17,0.58);
    }}
    .risk-axis-y {{
      left: -2px;
      top: 50%;
      transform: rotate(-90deg) translateY(-50%);
      transform-origin: left top;
    }}
    .risk-axis-x {{
      right: 18px;
      bottom: 6px;
    }}
    .risk-legend {{
      display: grid;
      gap: 8px;
    }}
    .risk-legend-item {{
      display: grid;
      grid-template-columns: 24px 1fr;
      gap: 10px;
      align-items: start;
      font-size: 12px;
      line-height: 1.3;
    }}
    .risk-legend-index {{
      width: 24px;
      height: 24px;
      border-radius: 50%;
      background: rgba(30,90,168,0.12);
      color: var(--accent);
      display: inline-flex;
      align-items: center;
      justify-content: center;
      font-weight: 700;
    }}
    .risk-item {{
      padding: 14px 16px;
      border-radius: 18px;
      border: 1px solid rgba(176, 43, 43, 0.18);
      background: rgba(176, 43, 43, 0.06);
      display: grid;
      gap: 8px;
    }}
    .risk-title {{
      font-size: 16px;
      font-weight: 700;
    }}
    .risk-meta,
    .risk-owner,
    .risk-mitigation {{
      font-size: 13px;
      color: rgba(17,17,17,0.78);
    }}
    .timeline-track {{
      display: grid;
      grid-template-columns: repeat(auto-fit, minmax(160px, 1fr));
      gap: 10px;
      position: relative;
    }}
    .timeline-track::before {{
      content: "";
      position: absolute;
      left: 10px;
      right: 10px;
      top: 20px;
      height: 2px;
      background: linear-gradient(90deg, rgba(30,90,168,0.2), rgba(30,90,168,0.5), rgba(30,90,168,0.2));
    }}
    .timeline-event {{
      padding: 14px 16px;
      border-radius: 18px;
      background: rgba(255,255,255,0.56);
      border: 1px solid var(--line);
      display: grid;
      gap: 8px;
      position: relative;
      z-index: 1;
    }}
    .timeline-timing {{
      font-size: 12px;
      text-transform: uppercase;
      letter-spacing: 0.12em;
      color: var(--accent);
    }}
    .timeline-label {{
      font-size: 15px;
      font-weight: 700;
    }}
    .timeline-dependency {{
      font-size: 13px;
      color: rgba(17,17,17,0.72);
    }}
    .comparison-table {{
      display: grid;
      gap: 8px;
    }}
    .chart-panel {{
      display: grid;
      gap: 10px;
      padding: 14px;
      border-radius: 18px;
      border: 1px solid rgba(30,90,168,0.12);
      background: rgba(255,255,255,0.54);
    }}
    .chart-header {{
      display: grid;
      gap: 4px;
    }}
    .chart-title {{
      font-size: 12px;
      text-transform: uppercase;
      letter-spacing: 0.08em;
      color: var(--accent);
    }}
    .chart-subtitle,
    .chart-footnote,
    .chart-fallback {{
      font-size: 12px;
      color: rgba(17,17,17,0.66);
    }}
    .echart-canvas {{
      height: 240px;
      width: 100%;
    }}
    .comparison-visual-grid {{
      display: grid;
      grid-template-columns: repeat(3, minmax(0, 1fr));
      gap: 12px;
    }}
    .comparison-visual-card {{
      padding: 14px;
      border-radius: 18px;
      background: rgba(255,255,255,0.56);
      border: 1px solid var(--line);
      display: grid;
      gap: 10px;
    }}
    .comparison-visual-dimension {{
      font-size: 11px;
      text-transform: uppercase;
      letter-spacing: 0.1em;
      color: var(--accent);
    }}
    .comparison-visual-lanes {{
      display: grid;
      gap: 8px;
    }}
    .comparison-lane {{
      display: grid;
      gap: 4px;
      padding: 10px;
      border-radius: 12px;
    }}
    .comparison-lane span {{
      font-size: 10px;
      text-transform: uppercase;
      letter-spacing: 0.08em;
      color: rgba(17,17,17,0.58);
    }}
    .comparison-lane strong {{
      font-size: 13px;
      line-height: 1.25;
    }}
    .comparison-lane-current {{
      background: rgba(17,17,17,0.05);
    }}
    .comparison-lane-target {{
      background: rgba(30,90,168,0.1);
    }}
    .comparison-visual-highlight {{
      font-size: 12px;
      line-height: 1.3;
      color: rgba(17,17,17,0.72);
    }}
    .metric-strip {{
      display: grid;
      grid-template-columns: repeat(auto-fit, minmax(180px, 1fr));
      gap: 12px;
    }}
    .metric-card {{
      padding: 18px;
      border-radius: 18px;
      background: linear-gradient(160deg, rgba(30,90,168,0.12), rgba(255,255,255,0.42));
      border: 1px solid rgba(30,90,168,0.12);
      display: grid;
      gap: 8px;
    }}
    .metric-card::after {{
      content: "";
      height: 8px;
      border-radius: 999px;
      background: linear-gradient(90deg, rgba(30,90,168,0.85), rgba(30,90,168,0.2));
    }}
    .metric-label {{
      font-size: 11px;
      text-transform: uppercase;
      letter-spacing: 0.12em;
      color: var(--accent);
    }}
    .metric-value {{
      font-size: 24px;
      font-weight: 700;
      line-height: 1;
    }}
    .metric-note {{
      font-size: 13px;
      color: rgba(17,17,17,0.72);
    }}
    .appendix-trace {{
      display: grid;
      gap: 10px;
    }}
    .appendix-trace-item {{
      padding: 16px 18px;
      border-radius: 16px;
      background: rgba(255,255,255,0.56);
      border: 1px dashed rgba(17,17,17,0.16);
      display: grid;
      gap: 6px;
    }}
    .appendix-trace-source {{
      font-size: 11px;
      text-transform: uppercase;
      letter-spacing: 0.12em;
      color: var(--accent);
    }}
    .appendix-trace-value {{
      font-size: 15px;
      font-weight: 600;
    }}
    .appendix-trace-note {{
      font-size: 13px;
      color: rgba(17,17,17,0.72);
    }}
    .comparison-row {{
      display: grid;
      grid-template-columns: 1.1fr 1fr 1fr 1fr;
      gap: 10px;
      padding: 12px 14px;
      border-radius: 16px;
      background: rgba(255,255,255,0.54);
      border: 1px solid var(--line);
      align-items: start;
      font-size: 13px;
    }}
    .comparison-header {{
      background: rgba(17,17,17,0.06);
      font-size: 12px;
      text-transform: uppercase;
      letter-spacing: 0.08em;
      color: rgba(17,17,17,0.72);
    }}
    .slide-footer {{
      display: grid;
      grid-template-columns: 1fr auto;
      gap: 12px;
      align-items: end;
    }}
    .visual-direction,
    .source-rail {{
      padding: 10px 12px;
      border-radius: 16px;
      background: rgba(17,17,17,0.04);
      font-size: 12px;
      line-height: 1.35;
    }}
    .composition-hero_statement h2 {{
      max-width: 15ch;
      font-size: 40px;
    }}
    .composition-hero_statement .core-message {{
      max-width: 38rem;
      font-size: 19px;
    }}
    .composition-decision_frame .core-message {{
      font-weight: 600;
    }}
    .composition-decision_frame .headline-block,
    .composition-risk_matrix .headline-block,
    .composition-timeline_with_dependencies .headline-block {{
      gap: 8px;
      max-width: min(72%, 820px);
    }}
    .composition-risk_matrix .risk-layout {{
      grid-template-columns: minmax(0, 0.92fr) minmax(0, 1.08fr);
      align-items: start;
    }}
    .composition-risk_matrix .risk-layout > :first-child {{
      grid-column: 1;
      grid-row: 1;
    }}
    .composition-risk_matrix .risk-layout > :nth-child(2) {{
      grid-column: 2;
      grid-row: 1 / span 2;
    }}
    .composition-risk_matrix .risk-layout > :last-child {{
      grid-column: 1;
      grid-row: 2;
    }}
    .slide-content.fit-compressed {{
      gap: 12px;
    }}
    .slide-content.fit-compressed h2 {{
      font-size: 28px;
    }}
    .slide-content.fit-compressed .core-message {{
      font-size: 15px;
    }}
    .slide-content.fit-compressed .evidence-card,
    .slide-content.fit-compressed .timeline-event,
    .slide-content.fit-compressed .decision-option,
    .slide-content.fit-compressed .risk-item,
    .slide-content.fit-compressed .comparison-row {{
      padding: 10px 12px;
    }}
    .slide-content.fit-compressed .hero-strip,
    .slide-content.fit-compressed .decision-ask {{
      padding: 12px 14px;
    }}
    .slide-content.fit-micro {{
      gap: 10px;
    }}
    .slide-content.fit-micro h2 {{
      font-size: 25px;
    }}
    .slide-content.fit-micro .core-message {{
      font-size: 14px;
    }}
    .slide-content.fit-micro .hero-claim,
    .slide-content.fit-micro .decision-option-summary,
    .slide-content.fit-micro .timeline-label,
    .slide-content.fit-micro .risk-title {{
      font-size: 14px;
    }}
    .slide-content.fit-micro .evidence-card,
    .slide-content.fit-micro .timeline-event,
    .slide-content.fit-micro .decision-option,
    .slide-content.fit-micro .risk-item,
    .slide-content.fit-micro .comparison-row {{
      padding: 9px 11px;
      gap: 7px;
    }}
    .slide-content.fit-micro .evidence-text,
    .slide-content.fit-micro .hero-summary,
    .slide-content.fit-micro .timeline-dependency,
    .slide-content.fit-micro .risk-meta,
    .slide-content.fit-micro .risk-owner,
    .slide-content.fit-micro .risk-mitigation {{
      font-size: 12px;
    }}
    .density-dense .slide-content,
    .composition-decision_frame .slide-content,
    .composition-risk_matrix .slide-content,
    .composition-timeline_with_dependencies .slide-content {{
      padding: 20px 22px 18px;
      gap: 12px;
    }}
    .density-dense h2,
    .composition-decision_frame h2,
    .composition-risk_matrix h2,
    .composition-timeline_with_dependencies h2 {{
      font-size: 29px;
      max-width: 20ch;
    }}
    .density-dense .core-message,
    .composition-decision_frame .core-message,
    .composition-risk_matrix .core-message,
    .composition-timeline_with_dependencies .core-message {{
      font-size: 16px;
    }}
    .slide-notes {{
      width: min(100%, var(--slide-stage-width));
      border-radius: 16px;
      background: rgba(255,255,255,0.42);
      border: 1px solid rgba(255,255,255,0.36);
      padding: 10px 14px;
    }}
    .slide-notes summary {{
      cursor: pointer;
      list-style: none;
      font-size: 12px;
      text-transform: uppercase;
      letter-spacing: 0.08em;
      color: var(--accent);
    }}
    .speaker-notes {{
      margin-top: 10px;
      font-size: 13px;
      line-height: 1.45;
      color: rgba(17,17,17,0.72);
    }}
    @media (max-width: 1440px), (max-height: 900px) {{
      .slide {{
        --slide-stage-width: min(1180px, calc(100vw - 40px));
      }}
      .slide-content {{
        padding: 18px 20px 16px;
        gap: 10px;
      }}
      h2 {{
        font-size: 28px;
      }}
      .core-message {{
        font-size: 15px;
      }}
      .composition-hero_statement h2 {{
        font-size: 34px;
      }}
      .composition-hero_statement .core-message {{
        font-size: 17px;
      }}
      .density-dense h2,
      .composition-decision_frame h2,
      .composition-risk_matrix h2,
      .composition-timeline_with_dependencies h2 {{
        font-size: 25px;
      }}
      .density-dense .core-message,
      .composition-decision_frame .core-message,
      .composition-risk_matrix .core-message,
      .composition-timeline_with_dependencies .core-message {{
        font-size: 14px;
      }}
      .evidence-card,
      .timeline-event,
      .decision-option,
      .risk-item,
      .comparison-row {{
        padding: 11px 12px;
      }}
      .hero-strip,
      .decision-ask {{
        padding: 12px 14px;
      }}
      .echart-canvas {{
        height: 210px;
      }}
    }}
    @media (max-width: 700px) {{
      .deck {{ padding: 12px; }}
      .slide {{
        --slide-stage-width: 100%;
      }}
      .slide-stage {{
        width: 100%;
      }}
      .slide-frame {{
        aspect-ratio: auto;
        min-height: auto;
        border-radius: 20px;
      }}
      .slide-content {{
        width: 100%;
        min-height: auto;
        transform: none;
        padding: 20px;
        grid-template-rows: auto;
      }}
      .headline-block {{
        max-width: 100%;
      }}
      .evidence-grid,
      .slide-footer,
      .hero-strip,
      .hero-layout,
      .map-body-layout,
      .risk-heatmap-panel,
      .comparison-visual-grid,
      .decision-layout,
      .comparison-row {{ grid-template-columns: 1fr; }}
      .matrix-map,
      .map-stage-row,
      .map-lane-row {{ grid-template-columns: 1fr; }}
    }}
  </style>
</head>
<body>
  <main class="deck">
    {''.join(slide_sections)}
  </main>
  <script>
    window.__PRODUCTOS_CHARTS__ = {json.dumps(chart_payloads)};
    window.__PRODUCTOS_fitSlides = function () {{
      var desktopViewport = window.matchMedia("(min-width: 701px)").matches;
      document.querySelectorAll(".slide-frame").forEach(function (frame) {{
        var content = frame.querySelector(".slide-content");
        if (!content) return;
        content.classList.remove("fit-compressed", "fit-micro");
        content.style.setProperty("--fit-scale", "1");
        if (!desktopViewport || frame.clientHeight === 0) return;

        var availableHeight = frame.clientHeight - 2;
        var naturalHeight = content.scrollHeight;
        if (naturalHeight <= availableHeight) return;

        content.classList.add("fit-compressed");
        naturalHeight = content.scrollHeight;
        if (naturalHeight > availableHeight * 1.04) {{
          content.classList.add("fit-micro");
          naturalHeight = content.scrollHeight;
        }}

        var scale = Math.min(1, availableHeight / naturalHeight);
        if (scale < 0.999) {{
          content.style.setProperty("--fit-scale", scale.toFixed(3));
        }}
      }});
    }};
    window.addEventListener("load", function () {{
      window.__PRODUCTOS_fitSlides();
      if (!window.echarts || !Array.isArray(window.__PRODUCTOS_CHARTS__)) return;
      window.__PRODUCTOS_CHARTS__.forEach(function (chartSpec) {{
        var element = document.getElementById(chartSpec.chart_id);
        if (!element) return;
        var chart = window.echarts.init(element);
        var colors = ["{theme['accent']}", "rgba(17,17,17,0.35)", "rgba(176,43,43,0.68)", "rgba(47,125,83,0.72)"];
        var option;
        if (chartSpec.chart_type === "heatmap") {{
          option = {{
            animationDuration: 500,
            tooltip: {{
              trigger: "item",
              formatter: function (params) {{
                return (params.data.label || "") + "<br/>Value: " + params.data.value[2];
              }}
            }},
            grid: {{ top: 28, bottom: 26 }},
            xAxis: {{ type: "category", data: chartSpec.x_categories }},
            yAxis: {{ type: "category", data: chartSpec.y_categories }},
            visualMap: {{
              min: 0,
              max: 100,
              orient: "horizontal",
              left: "center",
              bottom: 0,
              inRange: {{ color: ["#f4efe7", "{theme['accent']}"] }}
            }},
            series: [{{
              type: "heatmap",
              data: (chartSpec.heatmap_points || []).map(function (point) {{
                return {{ value: [point.x, point.y, point.value], label: point.label || "" }};
              }}),
              label: {{ show: true, formatter: function (params) {{ return params.data.label || params.data.value[2]; }} }},
              emphasis: {{ itemStyle: {{ shadowBlur: 10, shadowColor: "rgba(0,0,0,0.15)" }} }}
            }}]
          }};
        }} else {{
          var yAxes = (chartSpec.y_axes && chartSpec.y_axes.length ? chartSpec.y_axes : [chartSpec.unit || ""]).map(function (axisName, index) {{
            return {{
              type: "value",
              name: axisName,
              max: axisName === "count" ? null : 100,
              splitLine: index === 0 ? {{ lineStyle: {{ color: "rgba(17,17,17,0.08)" }} }} : {{ show: false }}
            }};
          }});
          option = {{
            animationDuration: 500,
            tooltip: {{ trigger: "axis" }},
            legend: {{ top: 0, textStyle: {{ fontFamily: "sans-serif" }} }},
            grid: {{ left: 40, right: 24, top: 34, bottom: 28 }},
            xAxis: {{
              type: "category",
              data: chartSpec.categories,
              axisLabel: {{ interval: 0, rotate: chartSpec.categories.length > 3 ? 20 : 0 }}
            }},
            yAxis: yAxes,
            series: (chartSpec.series || []).map(function (series, index) {{
              var effectiveType = series.type || (chartSpec.chart_type === "stacked_bar" ? "bar" : chartSpec.chart_type);
              return {{
                name: series.name,
                type: effectiveType === "combo" ? "line" : effectiveType,
                data: series.data,
                stack: series.stack || (chartSpec.chart_type === "stacked_bar" ? "total" : undefined),
                yAxisIndex: series.y_axis_index || 0,
                barMaxWidth: 28,
                smooth: effectiveType === "line",
                lineStyle: {{
                  type: series.name === "Baseline" || series.name === "Target" ? "dashed" : "solid"
                }},
                itemStyle: {{
                  color: colors[index % colors.length]
                }},
                emphasis: {{ focus: "series" }}
              }};
            }})
          }};
        }}
        chart.setOption(option);
        window.addEventListener("resize", function () {{
          window.__PRODUCTOS_fitSlides();
          chart.resize();
        }});
      }});
    }});
  </script>
</body>
</html>"""


def _ppt_shape_family_for_slide(slide: dict[str, Any]) -> str:
    composition_type = slide["composition_type"]
    if composition_type in _CANVAS_MAP_COMPOSITIONS:
        return "map_canvas"
    if composition_type in _MATRIX_MAP_COMPOSITIONS:
        return "matrix_map"
    return {
        "hero_statement": "hero_panel",
        "summary_cards": "summary_cards",
        "decision_frame": "decision_panel",
        "evidence_grid": "evidence_cards",
        "metric_strip": "metric_strip",
        "timeline_with_dependencies": "timeline_track",
        "risk_matrix": "risk_matrix",
        "comparison_table": "comparison_grid",
        "appendix_evidence": "appendix_trace",
    }[composition_type]


def _ppt_primitives_for_slide(slide: dict[str, Any]) -> list[str]:
    composition_type = slide["composition_type"]
    if composition_type in _CANVAS_MAP_COMPOSITIONS:
        return ["headline_panel", "stage_row", "lane_row", "map_cards", "evidence_rail"]
    if composition_type in _MATRIX_MAP_COMPOSITIONS:
        return ["headline_panel", "axis_labels", "quadrant_grid", "certainty_chips", "evidence_rail"]
    return {
        "hero_statement": ["headline_panel", "signal_rail", "recommendation_strip", "evidence_cards"],
        "summary_cards": ["headline_panel", "summary_chips", "evidence_cards"],
        "decision_frame": ["headline_panel", "decision_ask", "option_cards", "evidence_rail"],
        "evidence_grid": ["headline_panel", "evidence_cards"],
        "metric_strip": ["headline_panel", "metric_cards", "evidence_rail"],
        "timeline_with_dependencies": ["headline_panel", "timeline_track", "evidence_rail"],
        "risk_matrix": ["headline_panel", "risk_heatmap", "risk_cards", "evidence_rail"],
        "comparison_table": ["headline_panel", "comparison_grid", "evidence_rail"],
        "appendix_evidence": ["headline_panel", "appendix_trace", "evidence_rail"],
    }[composition_type]


def _ppt_instruction_for_slide(slide: dict[str, Any]) -> dict[str, Any]:
    composition_type = slide["composition_type"]
    map_payload = slide["composition_payload"].get("map_payload") or {}
    instruction = {
        "slide_id": slide["slide_id"],
        "composition_type": composition_type,
        "layout_variant": slide["layout_variant"],
        "native_shape_family": _ppt_shape_family_for_slide(slide),
        "fallback_layout": slide["ppt_render_hints"]["fallback_layout"],
        "prefer_native_shapes": slide["ppt_render_hints"]["prefer_native_shapes"],
        "target_profile": _slide_ppt_target_profile(slide),
        "primitives": _ppt_primitives_for_slide(slide),
        "evidence_mode": "cards" if slide["ppt_render_hints"].get("max_visible_evidence", 0) else "minimal",
        "connector_policy": (
            "directional_connectors"
            if composition_type in {"timeline_with_dependencies", "process_flow", "workflow_map", "roadmap_view"}
            else "grid_alignment"
        ),
        "source_refs": deepcopy(slide["source_refs"]),
    }
    if map_payload.get("map_type"):
        instruction["map_type"] = map_payload["map_type"]
    return instruction


def build_ppt_export_plan(render_spec: dict[str, Any]) -> dict[str, Any]:
    composition_counts: dict[str, int] = {}
    layout_counts: dict[str, int] = {}
    native_shape_counts: dict[str, int] = {}
    target_profile_counts: dict[str, int] = {}
    fallback_compositions = []
    slide_rendering_plan = []
    for slide in render_spec["slides"]:
        composition = slide["composition_type"]
        composition_counts[composition] = composition_counts.get(composition, 0) + 1
        layout = slide["ppt_render_hints"]["fallback_layout"]
        layout_counts[layout] = layout_counts.get(layout, 0) + 1
        native_shape_family = _ppt_shape_family_for_slide(slide)
        native_shape_counts[native_shape_family] = native_shape_counts.get(native_shape_family, 0) + 1
        target_profile = _slide_ppt_target_profile(slide)
        target_profile_counts[target_profile] = target_profile_counts.get(target_profile, 0) + 1
        slide_rendering_plan.append(_ppt_instruction_for_slide(slide))
        if target_profile == "ppt_safe":
            fallback_compositions.append(composition)

    notes = ["Generate native pptx from normalized render_spec compositions."]
    if composition_counts:
        notes.append(
            "Composition coverage: "
            + ", ".join(f"{name}={count}" for name, count in sorted(composition_counts.items()))
            + "."
        )
    if layout_counts:
        notes.append(
            "Layout coverage: "
            + ", ".join(f"{name}={count}" for name, count in sorted(layout_counts.items()))
            + "."
        )
    if fallback_compositions:
        notes.append(
            "PPT-safe fallback rendering is planned for: " + ", ".join(sorted(set(fallback_compositions))) + "."
        )
    if native_shape_counts:
        notes.append(
            "Native shape families: "
            + ", ".join(f"{name}={count}" for name, count in sorted(native_shape_counts.items()))
            + "."
        )
    if target_profile_counts:
        notes.append(
            "Target profiles: "
            + ", ".join(f"{name}={count}" for name, count in sorted(target_profile_counts.items()))
            + "."
        )

    return {
        "schema_version": "1.0.0",
        "ppt_export_plan_id": f"ppt_export_plan_{render_spec['render_spec_id']}",
        "render_spec_id": render_spec["render_spec_id"],
        "workspace_id": render_spec["workspace_id"],
        "engine": "pptxgenjs",
        "status": "composition-aware",
        "aspect_ratio": render_spec["aspect_ratio"],
        "theme_preset": render_spec["theme"]["preset"],
        "slide_count": len(render_spec["slides"]),
        "composition_counts": composition_counts,
        "layout_counts": layout_counts,
        "native_shape_counts": native_shape_counts,
        "slide_rendering_plan": slide_rendering_plan,
        "fidelity_status": "planned_fallbacks" if "ppt_safe" in target_profile_counts else "aligned",
        "notes": notes,
        "generated_at": now_iso(),
    }


def _hex_to_rgb(color: str) -> tuple[int, int, int]:
    normalized = color.lstrip("#")
    if len(normalized) != 6:
        return (30, 90, 168)
    return tuple(int(normalized[index:index + 2], 16) for index in (0, 2, 4))


def _ppt_apply_shape_style(shape: Any, *, fill_rgb: tuple[int, int, int], line_rgb: tuple[int, int, int]) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shape.line.color.rgb = RGBColor(*line_rgb)


def _ppt_add_textbox(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color_rgb: tuple[int, int, int] = (17, 17, 17),
    align: Any = None,
) -> Any:
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color_rgb)
    return textbox


def _ppt_add_panel(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str | None = None,
    body: str | None = None,
    *,
    fill_rgb: tuple[int, int, int] = (248, 246, 241),
    line_rgb: tuple[int, int, int] = (210, 217, 226),
) -> Any:
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    _ppt_apply_shape_style(panel, fill_rgb=fill_rgb, line_rgb=line_rgb)
    if title:
        _ppt_add_textbox(slide, left + 0.15, top + 0.08, width - 0.3, 0.28, title, font_size=11, bold=True, color_rgb=line_rgb)
    if body:
        _ppt_add_textbox(slide, left + 0.15, top + 0.34, width - 0.3, height - 0.42, body, font_size=12, color_rgb=(17, 17, 17))
    return panel


def _ppt_add_card_grid(
    slide: Any,
    items: list[dict[str, Any]],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    columns: int = 2,
    fill_rgb: tuple[int, int, int] = (248, 246, 241),
    line_rgb: tuple[int, int, int] = (210, 217, 226),
) -> None:
    if not items:
        return
    visible = items[: min(len(items), 6)]
    rows = max(1, (len(visible) + columns - 1) // columns)
    gap_x = 0.16
    gap_y = 0.16
    card_width = (width - gap_x * (columns - 1)) / columns
    card_height = (height - gap_y * (rows - 1)) / rows
    for index, item in enumerate(visible):
        row = index // columns
        col = index % columns
        card_left = left + col * (card_width + gap_x)
        card_top = top + row * (card_height + gap_y)
        title = item.get("label") or item.get("dimension") or item.get("timing") or "Item"
        body_parts = [
            item.get("value"),
            item.get("description"),
            item.get("annotation"),
            item.get("highlight"),
            item.get("mitigation"),
        ]
        body = "\n".join(part for part in body_parts if part)
        _ppt_add_panel(
            slide,
            card_left,
            card_top,
            card_width,
            card_height,
            title,
            body,
            fill_rgb=fill_rgb,
            line_rgb=line_rgb,
        )


def _ppt_add_header(slide: Any, slide_data: dict[str, Any], accent_rgb: tuple[int, int, int]) -> None:
    _ppt_add_textbox(slide, 0.45, 0.2, 9.2, 0.45, slide_data["title"], font_size=26, bold=True)
    _ppt_add_textbox(slide, 0.45, 0.75, 9.8, 0.45, slide_data["headline"], font_size=16, color_rgb=accent_rgb)
    _ppt_add_textbox(slide, 10.75, 0.25, 2.0, 0.25, slide_data["confidence_state"].upper(), font_size=10, bold=True, color_rgb=accent_rgb, align=PP_ALIGN.RIGHT)


def _ppt_render_hero_panel(slide: Any, slide_data: dict[str, Any], accent_rgb: tuple[int, int, int]) -> None:
    payload = slide_data["composition_payload"]
    _ppt_add_panel(slide, 0.45, 1.25, 7.1, 2.0, "Primary claim", payload.get("primary_claim", slide_data["core_message"]), fill_rgb=(244, 248, 252), line_rgb=accent_rgb)
    evidence = payload.get("evidence_items", [])
    _ppt_add_card_grid(slide, evidence, left=7.8, top=1.25, width=4.85, height=2.8, columns=1, fill_rgb=(249, 247, 242), line_rgb=(210, 217, 226))
    recommendation = payload.get("recommendation") or "No recommendation captured."
    decision_ask = payload.get("decision_ask") or "No explicit ask."
    strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(5.5), Inches(12.2), Inches(1.0))
    _ppt_apply_shape_style(strip, fill_rgb=(236, 242, 250), line_rgb=accent_rgb)
    _ppt_add_textbox(slide, 0.7, 5.72, 5.7, 0.25, "Recommendation", font_size=10, bold=True, color_rgb=accent_rgb)
    _ppt_add_textbox(slide, 0.7, 5.96, 5.4, 0.38, recommendation, font_size=13)
    _ppt_add_textbox(slide, 6.7, 5.72, 5.4, 0.25, "Ask", font_size=10, bold=True, color_rgb=accent_rgb)
    _ppt_add_textbox(slide, 6.7, 5.96, 5.1, 0.38, decision_ask, font_size=13)


def _ppt_render_map_canvas(slide: Any, slide_data: dict[str, Any], accent_rgb: tuple[int, int, int]) -> None:
    map_payload = slide_data["composition_payload"].get("map_payload") or {}
    body_left = 0.45
    body_width = 12.2
    if map_payload.get("stages"):
        stage_width = max(1.3, body_width / max(1, len(map_payload["stages"][:6])) - 0.08)
        for index, stage in enumerate(map_payload["stages"][:6]):
            _ppt_add_panel(slide, body_left + index * (stage_width + 0.08), 1.35, stage_width, 0.42, stage, None, fill_rgb=(236, 242, 250), line_rgb=accent_rgb)
    if map_payload.get("lanes"):
        lane_width = max(1.3, body_width / max(1, len(map_payload["lanes"][:6])) - 0.08)
        for index, lane in enumerate(map_payload["lanes"][:6]):
            _ppt_add_panel(slide, body_left + index * (lane_width + 0.08), 1.9, lane_width, 0.38, lane, None, fill_rgb=(243, 243, 243), line_rgb=(150, 150, 150))
    items = map_payload.get("items") or map_payload.get("nodes") or []
    _ppt_add_card_grid(slide, items, left=0.45, top=2.45, width=8.0, height=3.55, columns=2, fill_rgb=(249, 247, 242), line_rgb=(210, 217, 226))
    evidence = slide_data["composition_payload"].get("evidence_items", [])
    _ppt_add_card_grid(slide, evidence, left=8.7, top=2.45, width=3.95, height=3.55, columns=1, fill_rgb=(244, 248, 252), line_rgb=accent_rgb)
    if map_payload.get("axes"):
        _ppt_add_textbox(slide, 8.7, 1.95, 3.95, 0.22, "Axes: " + " / ".join(map_payload["axes"][:4]), font_size=10, color_rgb=accent_rgb)
    if map_payload.get("certainty_notes"):
        _ppt_add_textbox(slide, 0.45, 6.2, 12.0, 0.28, "Certainty: " + " | ".join(map_payload["certainty_notes"][:3]), font_size=10, color_rgb=accent_rgb)


def _ppt_render_matrix_map(slide: Any, slide_data: dict[str, Any], accent_rgb: tuple[int, int, int]) -> None:
    map_payload = slide_data["composition_payload"].get("map_payload") or {}
    axes = map_payload.get("axes", [])
    items = map_payload.get("items", [])[:4]
    cell_left = 0.9
    cell_top = 1.6
    cell_width = 4.0
    cell_height = 1.7
    gap = 0.2
    _ppt_add_textbox(slide, 0.25, 3.15, 0.35, 1.2, axes[0] if len(axes) > 0 else "Axis 1", font_size=10, bold=True, color_rgb=accent_rgb)
    _ppt_add_textbox(slide, 5.1, 5.25, 2.0, 0.25, axes[1] if len(axes) > 1 else "Axis 2", font_size=10, bold=True, color_rgb=accent_rgb)
    for index in range(4):
        row = index // 2
        col = index % 2
        item = items[index] if index < len(items) else {"label": "Open quadrant", "description": "No mapped item yet."}
        _ppt_add_panel(
            slide,
            cell_left + col * (cell_width + gap),
            cell_top + row * (cell_height + gap),
            cell_width,
            cell_height,
            item.get("label", "Quadrant"),
            item.get("description"),
            fill_rgb=(249, 247, 242),
            line_rgb=accent_rgb if index < len(items) else (210, 217, 226),
        )
    evidence = slide_data["composition_payload"].get("evidence_items", [])
    _ppt_add_card_grid(slide, evidence, left=9.15, top=1.8, width=3.5, height=3.8, columns=1, fill_rgb=(244, 248, 252), line_rgb=accent_rgb)
    if map_payload.get("certainty_notes"):
        _ppt_add_textbox(slide, 0.9, 5.8, 7.8, 0.28, "Certainty: " + " | ".join(map_payload["certainty_notes"][:3]), font_size=10, color_rgb=accent_rgb)


def _ppt_render_timeline(slide: Any, slide_data: dict[str, Any], accent_rgb: tuple[int, int, int]) -> None:
    events = slide_data["composition_payload"].get("timeline_events", [])[:4]
    if not events:
        return
    start_left = 0.85
    top = 2.3
    width = 11.4
    event_width = max(1.8, width / len(events) - 0.18)
    line_y = top + 0.95
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(start_left), Inches(line_y), Inches(start_left + width), Inches(line_y))
    for index, event in enumerate(events):
        left = start_left + index * (event_width + 0.18)
        slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.15), Inches(line_y - 0.1), Inches(0.18), Inches(0.18))
        _ppt_add_panel(slide, left, top + 0.25, event_width, 1.5, event.get("timing"), event.get("label"), fill_rgb=(249, 247, 242), line_rgb=accent_rgb)
        if event.get("dependency"):
            _ppt_add_textbox(slide, left, top + 1.8, event_width, 0.35, event["dependency"], font_size=10, color_rgb=(90, 90, 90))


def _ppt_render_risk_matrix(slide: Any, slide_data: dict[str, Any], accent_rgb: tuple[int, int, int]) -> None:
    risks = slide_data["composition_payload"].get("risk_items", [])[:3]
    grid_left = 0.75
    grid_top = 1.7
    cell = 0.85
    for row in range(3):
        for col in range(3):
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(grid_left + col * cell), Inches(grid_top + row * cell), Inches(cell - 0.05), Inches(cell - 0.05))
            _ppt_apply_shape_style(rect, fill_rgb=(248 - row * 10, 244 - col * 6, 236), line_rgb=(220, 220, 220))
    for index, risk in enumerate(risks):
        left = grid_left + min(index, 2) * cell + 0.25
        top = grid_top + max(0, 2 - index) * cell + 0.2
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(0.28), Inches(0.28))
        _ppt_apply_shape_style(marker, fill_rgb=accent_rgb, line_rgb=accent_rgb)
        _ppt_add_textbox(slide, left + 0.45, top - 0.05, 3.5, 0.6, risk.get("label", "Risk"), font_size=11)
    evidence = slide_data["composition_payload"].get("evidence_items", [])
    _ppt_add_card_grid(slide, evidence, left=4.0, top=1.8, width=8.65, height=3.8, columns=2, fill_rgb=(249, 247, 242), line_rgb=(210, 217, 226))


def _ppt_render_generic_cards(slide: Any, slide_data: dict[str, Any], accent_rgb: tuple[int, int, int]) -> None:
    payload = slide_data["composition_payload"]
    items = payload.get("evidence_items") or payload.get("comparison_rows") or []
    if slide_data["composition_type"] == "decision_frame":
        options = payload.get("options", [])
        items = [
            {
                "label": option["label"],
                "value": option["summary"],
                "annotation": " | ".join(option.get("tradeoffs", [])[:2]),
            }
            for option in options
        ] or items
    if slide_data["composition_type"] == "metric_strip":
        items = [
            {
                "label": row["dimension"],
                "value": row["target_state"],
                "annotation": row["highlight"],
            }
            for row in payload.get("comparison_rows", [])
        ]
    _ppt_add_card_grid(slide, items, left=0.55, top=1.7, width=12.0, height=4.6, columns=2, fill_rgb=(249, 247, 242), line_rgb=(210, 217, 226))
    _ppt_add_textbox(slide, 0.55, 6.35, 11.8, 0.4, slide_data["core_message"], font_size=12, color_rgb=accent_rgb)


def _ppt_render_slide(slide: Any, slide_data: dict[str, Any], instruction: dict[str, Any], render_spec: dict[str, Any]) -> None:
    accent_rgb = _hex_to_rgb(render_spec["theme"]["accent"])
    _ppt_add_header(slide, slide_data, accent_rgb)
    family = instruction["native_shape_family"]
    if family == "hero_panel":
        _ppt_render_hero_panel(slide, slide_data, accent_rgb)
        return
    if family == "map_canvas":
        _ppt_render_map_canvas(slide, slide_data, accent_rgb)
        return
    if family == "matrix_map":
        _ppt_render_matrix_map(slide, slide_data, accent_rgb)
        return
    if family == "timeline_track":
        _ppt_render_timeline(slide, slide_data, accent_rgb)
        return
    if family == "risk_matrix":
        _ppt_render_risk_matrix(slide, slide_data, accent_rgb)
        return
    _ppt_render_generic_cards(slide, slide_data, accent_rgb)


def write_ppt_presentation(render_spec: dict[str, Any], output_path: str | Path) -> Path:
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed, so native PPT export is unavailable.")

    export_plan = build_ppt_export_plan(render_spec)
    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    if render_spec["aspect_ratio"] == "4:3":
        presentation.slide_width = Inches(10)
        presentation.slide_height = Inches(7.5)
    else:
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
    presentation.core_properties.author = "ProductOS"
    presentation.core_properties.title = render_spec["render_spec_id"]
    presentation.core_properties.subject = export_plan["ppt_export_plan_id"]

    instructions = {item["slide_id"]: item for item in export_plan["slide_rendering_plan"]}
    blank_layout = presentation.slide_layouts[6]
    for slide_data in render_spec["slides"]:
        slide = presentation.slides.add_slide(blank_layout)
        instruction = instructions[slide_data["slide_id"]]
        _ppt_render_slide(slide, slide_data, instruction, render_spec)

    presentation.save(str(path))
    return path


def _presentation_mode_for_rendering_mode(rendering_mode: str) -> str:
    return presentation_mode_for_map_rendering_mode(rendering_mode)


def _map_fallback_layout(map_type: str) -> str:
    return fallback_layout_for_composition(map_type)


def _slide_spec_layout_for_composition(composition_type: str) -> str:
    return {
        "hero_statement": "cover",
        "summary_cards": "headline_grid",
        "decision_frame": "two_column",
        "evidence_grid": "headline_grid",
        "metric_strip": "two_column",
        "timeline_with_dependencies": "timeline",
        "risk_matrix": "risk_matrix",
        "comparison_table": "headline_grid",
        "appendix_evidence": "closing",
        "roadmap_view": "map_canvas",
        "user_journey_map": "map_canvas",
        "process_flow": "map_canvas",
        "workflow_map": "map_canvas",
        "capability_map": "map_canvas",
        "product_map": "map_canvas",
        "feature_map": "map_canvas",
        "mind_map": "map_canvas",
        "swot_matrix": "matrix_map",
        "impact_effort_matrix": "matrix_map",
    }[composition_type]


def _slide_spec_block_type_for_composition(composition_type: str) -> str:
    return {
        "hero_statement": "headline",
        "summary_cards": "bullets",
        "decision_frame": "bullets",
        "evidence_grid": "bullets",
        "metric_strip": "stat",
        "timeline_with_dependencies": "timeline_item",
        "risk_matrix": "risk",
        "comparison_table": "quote",
        "appendix_evidence": "closing_note",
        "roadmap_view": "map_canvas",
        "user_journey_map": "map_canvas",
        "process_flow": "map_canvas",
        "workflow_map": "map_canvas",
        "capability_map": "map_canvas",
        "product_map": "map_canvas",
        "feature_map": "map_canvas",
        "mind_map": "map_canvas",
        "swot_matrix": "matrix_map",
        "impact_effort_matrix": "matrix_map",
    }[composition_type]


def _slide_spec_block_data(slide: dict[str, Any]) -> dict[str, Any] | None:
    if slide["composition_type"] not in _MAP_COMPOSITIONS:
        return None
    map_payload = slide["composition_payload"].get("map_payload") or {}
    return {
        "map_type": map_payload.get("map_type", slide["composition_type"]),
        "decision_use_case": map_payload.get("decision_use_case", slide["core_message"]),
        "certainty_notes": deepcopy(map_payload.get("certainty_notes", [])),
        "stages": deepcopy(map_payload.get("stages", [])),
        "lanes": deepcopy(map_payload.get("lanes", [])),
        "axes": deepcopy(map_payload.get("axes", [])),
        "items": deepcopy(map_payload.get("items", [])),
        "nodes": deepcopy(map_payload.get("nodes", [])),
    }


def _slide_spec_block_content(slide: dict[str, Any]) -> list[str] | str:
    if slide["composition_type"] in _MAP_COMPOSITIONS:
        map_payload = slide["composition_payload"].get("map_payload") or {}
        items = map_payload.get("items") or map_payload.get("nodes") or []
        labels = [item["label"] for item in items[:4] if item.get("label")]
        return labels or [slide["core_message"]]
    if slide["composition_type"] == "hero_statement":
        return slide["headline"]
    return [item["value"] for item in _payload_evidence_items(slide)] or [slide["core_message"]]


def build_visual_map_render_spec(
    visual_map_spec: dict[str, Any],
    *,
    theme_name: str = "atlas",
    aspect_ratio: str = "16:9",
) -> dict[str, Any]:
    presentation_mode = _presentation_mode_for_rendering_mode(visual_map_spec["rendering_mode"])
    presentation_format = presentation_format_for_map_rendering_mode(visual_map_spec["rendering_mode"])
    synthetic_brief = {
        "presentation_mode": presentation_mode,
        "density_preference": density_preference_for_map_rendering_mode(visual_map_spec["rendering_mode"]),
        "presentation_format": presentation_format,
        "theme_preset": theme_name,
        "customer_safe": False,
        "redaction_policy": "internal_only",
    }
    outline = {"appendix_link": False, "must_show_owner": False, "must_show_risk": False}
    map_payload = {
        "map_type": visual_map_spec["map_type"],
        "decision_use_case": visual_map_spec["decision_use_case"],
        "certainty_notes": deepcopy(visual_map_spec.get("certainty_notes", [])),
        "nodes": deepcopy(visual_map_spec["payload"].get("nodes", [])),
        "lanes": deepcopy(visual_map_spec["payload"].get("lanes", [])),
        "stages": deepcopy(visual_map_spec["payload"].get("stages", [])),
        "axes": deepcopy(visual_map_spec["payload"].get("axes", [])),
        "items": deepcopy(visual_map_spec["payload"].get("items", [])),
    }
    evidence_items = [
        {
            "evidence_id": f"evidence_{visual_map_spec['visual_map_spec_id']}_{index + 1}",
            "label": f"Source {index + 1}",
            "evidence_type": "source_artifact",
            "value": source_ref,
            "source_ref": source_ref,
            "annotation": "Linked source artifact used to construct this map.",
        }
        for index, source_ref in enumerate(visual_map_spec["source_artifact_ids"][:4])
    ]
    if visual_map_spec.get("certainty_notes") and evidence_items:
        evidence_items[0]["annotation"] = visual_map_spec["certainty_notes"][0]

    slide = _apply_render_budgets(
        {
            "slide_id": visual_map_spec["visual_map_spec_id"],
            "title": visual_map_spec["title"],
            "composition_type": visual_map_spec["map_type"],
            "headline": visual_map_spec["title"],
            "core_message": visual_map_spec["primary_message"],
            "composition_payload": {
                "primary_claim": visual_map_spec["primary_message"],
                "summary": visual_map_spec["decision_use_case"],
                "evidence_items": evidence_items,
                "map_payload": map_payload,
            },
            "layout_variant": _map_layout_variant(visual_map_spec["map_type"]),
            "density_mode": _density_mode(synthetic_brief),
            "visual_tokens": _visual_tokens(synthetic_brief, visual_map_spec["map_type"]),
            "visibility_rules": _visibility_rules(synthetic_brief, outline),
            "annotation_rules": _annotation_rules(outline),
            "html_render_hints": {
                "prefer_native_shapes": True,
                "fallback_layout": _map_fallback_layout(visual_map_spec["map_type"]),
                "target_profile": _html_target_profile(synthetic_brief, visual_map_spec["map_type"]),
                "emphasize_claim": True,
                "show_evidence_as_cards": True,
            },
            "ppt_render_hints": {
                "prefer_native_shapes": True,
                "fallback_layout": _map_fallback_layout(visual_map_spec["map_type"]),
                "target_profile": _ppt_target_profile(
                    synthetic_brief,
                    visual_map_spec["map_type"],
                    _html_target_profile(synthetic_brief, visual_map_spec["map_type"]),
                ),
            },
            "source_refs": deepcopy(visual_map_spec["source_artifact_ids"]),
            "speaker_notes": (
                f"Use this map to support {visual_map_spec['decision_use_case']} "
                f"The main takeaway is: {visual_map_spec['primary_message']}"
            ),
            "visual_direction": (
                "Render a structured PM map with explicit labels, evidence chips, and visible uncertainty."
            ),
            "confidence_state": "medium",
            "variant_rules": {
                "live": "Keep the visual map dominant and reduce evidence to the essential labels.",
                "async": "Keep the visual map and evidence cards visible together.",
                "memo": "Retain the map and add enough evidence context to explain the takeaway.",
            },
        },
        synthetic_brief,
    )

    return {
        "schema_version": "1.0.0",
        "render_spec_id": f"render_spec_{visual_map_spec['visual_map_spec_id']}",
        "presentation_story_id": f"presentation_story_{visual_map_spec['visual_map_spec_id']}",
        "presentation_brief_id": visual_map_spec["visual_map_spec_id"],
        "workspace_id": visual_map_spec["workspace_id"],
        "presentation_mode": presentation_mode,
        "theme": theme_preset(theme_name),
        "aspect_ratio": aspect_ratio,
        "slides": [slide],
        "generated_at": now_iso(),
    }


def render_visual_map_spec_html(
    visual_map_spec: dict[str, Any],
    *,
    theme_name: str = "atlas",
    aspect_ratio: str = "16:9",
) -> str:
    return render_render_spec_html(
        build_visual_map_render_spec(
            visual_map_spec,
            theme_name=theme_name,
            aspect_ratio=aspect_ratio,
        )
    )


def write_html_presentation(render_spec: dict[str, Any], output_path: str | Path) -> Path:
    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(render_render_spec_html(render_spec), encoding="utf-8")
    return path


def write_render_spec_payload(render_spec: dict[str, Any], output_path: str | Path) -> Path:
    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(render_spec, indent=2), encoding="utf-8")
    return path


def build_slide_spec(
    presentation_brief: dict[str, Any],
    aspect_ratio: str = "16:9",
) -> dict[str, Any]:
    render_spec = build_render_spec(
        presentation_brief,
        build_presentation_story(presentation_brief, build_evidence_pack(presentation_brief)),
        aspect_ratio=aspect_ratio,
    )
    return {
        "schema_version": "1.0.0",
        "slide_spec_id": f"slide_spec_{presentation_brief['presentation_brief_id']}",
        "presentation_brief_id": presentation_brief["presentation_brief_id"],
        "workspace_id": presentation_brief["workspace_id"],
        "theme": render_spec["theme"],
        "aspect_ratio": render_spec["aspect_ratio"],
        "slides": [
            {
                "slide_id": slide["slide_id"],
                "title": slide["title"],
                "layout": _slide_spec_layout_for_composition(slide["composition_type"]),
                "eyebrow": presentation_brief["audience"],
                "blocks": [
                    {
                        **{
                            "block_type": _slide_spec_block_type_for_composition(slide["composition_type"]),
                            "content": _slide_spec_block_content(slide),
                        },
                        **({"data": _slide_spec_block_data(slide)} if _slide_spec_block_data(slide) else {}),
                    }
                ],
            }
            for slide in render_spec["slides"]
        ],
        "generated_at": now_iso(),
    }


def build_visual_map_slide_spec(
    visual_map_spec: dict[str, Any],
    *,
    theme_name: str = "atlas",
    aspect_ratio: str = "16:9",
) -> dict[str, Any]:
    render_spec = build_visual_map_render_spec(
        visual_map_spec,
        theme_name=theme_name,
        aspect_ratio=aspect_ratio,
    )
    slide = render_spec["slides"][0]
    return {
        "schema_version": "1.0.0",
        "slide_spec_id": f"slide_spec_{visual_map_spec['visual_map_spec_id']}",
        "presentation_brief_id": visual_map_spec["visual_map_spec_id"],
        "workspace_id": visual_map_spec["workspace_id"],
        "theme": render_spec["theme"],
        "aspect_ratio": render_spec["aspect_ratio"],
        "slides": [
            {
                "slide_id": slide["slide_id"],
                "title": slide["title"],
                "layout": _slide_spec_layout_for_composition(slide["composition_type"]),
                "eyebrow": visual_map_spec["audience"],
                "blocks": [
                    {
                        "block_type": _slide_spec_block_type_for_composition(slide["composition_type"]),
                        "content": _slide_spec_block_content(slide),
                        "data": _slide_spec_block_data(slide),
                    }
                ],
            }
        ],
        "generated_at": now_iso(),
    }


def render_slide_spec_html(slide_spec: dict[str, Any]) -> str:
    render_spec = {
        "render_spec_id": slide_spec["slide_spec_id"],
        "presentation_mode": "async",
        "theme": slide_spec["theme"],
        "slides": [
            {
                "slide_id": slide["slide_id"],
                "title": slide["title"],
                "composition_type": (
                    slide["blocks"][0].get("data", {}).get("map_type")
                    or {
                        "cover": "hero_statement",
                        "headline_grid": "summary_cards",
                        "two_column": "summary_cards",
                        "timeline": "timeline_with_dependencies",
                        "risk_matrix": "risk_matrix",
                        "closing": "appendix_evidence",
                        "map_canvas": "user_journey_map",
                        "matrix_map": "impact_effort_matrix",
                    }[slide["layout"]]
                ),
                "headline": slide["title"],
                "core_message": slide["title"],
                "composition_payload": {
                    "primary_claim": slide["title"],
                    "evidence_items": [
                        {
                            "evidence_id": f"{slide['slide_id']}_legacy_{index + 1}",
                            "label": "Legacy support",
                            "evidence_type": "fact",
                            "value": item,
                            "source_ref": slide.get("eyebrow", "source_unavailable"),
                            "annotation": "Legacy slide-spec compatibility mode.",
                        }
                        for index, item in enumerate(
                            slide["blocks"][0]["content"]
                            if isinstance(slide["blocks"][0]["content"], list)
                            else [str(slide["blocks"][0]["content"])]
                        )
                    ],
                    "summary": slide["title"],
                    **(
                        {"map_payload": slide["blocks"][0]["data"]}
                        if slide["blocks"][0].get("data")
                        else {}
                    ),
                },
                "layout_variant": slide["layout"],
                "density_mode": "balanced",
                "visual_tokens": {
                    "theme_preset": slide_spec["theme"]["preset"],
                    "accent": slide_spec["theme"]["accent"],
                    "surface_style": "layered_panel",
                    "emphasis_style": "evidence_support",
                },
                "visibility_rules": {
                    "show_sources": True,
                    "show_confidence": True,
                    "show_appendix_link": False,
                    "customer_safe": False,
                    "redaction_policy": "none",
                },
                "annotation_rules": {
                    "show_confidence_badge": True,
                    "show_owner_labels": False,
                    "show_risk_labels": False,
                    "show_source_tags": True,
                },
                "html_render_hints": {
                    "prefer_native_shapes": True,
                    "fallback_layout": slide["layout"],
                    "target_profile": "dual_target",
                    "emphasize_claim": True,
                    "show_evidence_as_cards": True,
                },
                "ppt_render_hints": {
                    "prefer_native_shapes": True,
                    "fallback_layout": slide["layout"],
                    "target_profile": "dual_target",
                },
                "source_refs": [slide.get("eyebrow", "source_unavailable")],
                "speaker_notes": "Legacy slide spec converted for HTML rendering.",
                "visual_direction": "Legacy slide-spec compatibility mode.",
                "confidence_state": "medium",
                "variant_rules": {
                    "live": "Condense to headline only.",
                    "async": "Keep supporting bullets visible.",
                    "memo": "Expand bullets into prose.",
                },
            }
            for slide in slide_spec["slides"]
        ],
    }
    return render_render_spec_html(render_spec)


def write_slide_spec_payload(slide_spec: dict[str, Any], output_path: str | Path) -> Path:
    return write_render_spec_payload(slide_spec, output_path)
